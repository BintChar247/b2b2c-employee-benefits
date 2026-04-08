"""
MUFG Employee Benefits Portal — English Web Presentation
Finetuned with B2B2C Global Market Review insights.
Clean, data-led, web-shareable design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand Palette ─────────────────────────────────────────────────────────────
RED     = RGBColor(0xC8, 0x10, 0x2E)   # MUFG red
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
GOLD    = RGBColor(0xB8, 0x86, 0x0B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGREY   = RGBColor(0xF5, 0xF5, 0xF5)
MGREY   = RGBColor(0xAA, 0xAA, 0xAA)
DGREY   = RGBColor(0x44, 0x44, 0x44)
BLUE    = RGBColor(0x00, 0x5F, 0xAF)
GREEN   = RGBColor(0x1A, 0x7A, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Primitives ────────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h,
        size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def para(tf, text, size=12, bold=False, color=DGREY,
         align=PP_ALIGN.LEFT, before=Pt(5), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = before
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


# ── Shared chrome ─────────────────────────────────────────────────────────────

def header(slide, title, sub=""):
    rect(slide, 0, 0, W, H, fill=WHITE)
    rect(slide, 0, 0, W, Inches(1.05), fill=DARK)
    rect(slide, 0, Inches(1.05), Inches(0.06), H - Inches(1.05), fill=RED)
    # MUFG badge
    rect(slide, W - Inches(1.75), Inches(0.13), Inches(1.55), Inches(0.72), fill=RED)
    txt(slide, "MUFG", W - Inches(1.75), Inches(0.18), Inches(1.55), Inches(0.55),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, Inches(0.28), Inches(0.12), Inches(10.5), Inches(0.62),
        size=22, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.28), Inches(0.7), Inches(10.5), Inches(0.32),
            size=11, color=MGREY, italic=True)
    # footer
    rect(slide, 0, H - Inches(0.32), W, Inches(0.32), fill=LGREY)
    txt(slide, "CONFIDENTIAL  ·  MUFG Employee Benefits Portal  ·  Internal Proposal  ·  April 2026",
        Inches(0.28), H - Inches(0.30), W - Inches(0.5), Inches(0.26),
        size=8, color=MGREY)


def divider(prs, num, title, sub=""):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, W, H, fill=DARK)
    rect(sl, 0, 0, Inches(0.55), H, fill=RED)
    rect(sl, 0, H - Inches(0.08), W, Inches(0.08), fill=RED)
    txt(sl, f"{num:02d}", Inches(1.0), Inches(1.8), Inches(2), Inches(1.6),
        size=80, bold=True, color=RED, align=PP_ALIGN.LEFT)
    txt(sl, title, Inches(1.0), Inches(3.4), Inches(11), Inches(1.1),
        size=36, bold=True, color=WHITE)
    if sub:
        rect(sl, Inches(1.0), Inches(4.6), Inches(3.5), Inches(0.04), fill=GOLD)
        txt(sl, sub, Inches(1.0), Inches(4.75), Inches(11), Inches(0.55),
            size=15, color=MGREY, italic=True)
    txt(sl, "MUFG  ·  CONFIDENTIAL", W - Inches(3), H - Inches(0.45), Inches(2.8), Inches(0.38),
        size=9, color=MGREY, align=PP_ALIGN.RIGHT)
    return sl


def stat_box(slide, num, label, bx, by, w=Inches(2.9), h=Inches(1.85),
             bg=WHITE, border=RED, num_color=RED):
    rect(slide, bx, by, w, h, fill=bg, line=border, lw=Pt(1.5))
    rect(slide, bx, by, w, Inches(0.06), fill=border)
    txt(slide, num,   bx + Inches(0.15), by + Inches(0.18), w - Inches(0.3), Inches(0.9),
        size=32, bold=True, color=num_color, align=PP_ALIGN.LEFT)
    txt(slide, label, bx + Inches(0.15), by + Inches(1.05), w - Inches(0.3), Inches(0.65),
        size=11, color=DGREY)


def card(slide, title, bullets, bx, by, w, h,
         title_bg=RED, title_color=WHITE, bullet_color=DARK, bullet_size=11):
    rect(slide, bx, by, w, h, fill=WHITE, line=title_bg, lw=Pt(1))
    rect(slide, bx, by, w, Inches(0.48), fill=title_bg)
    txt(slide, title, bx + Inches(0.12), by + Inches(0.07), w - Inches(0.24), Inches(0.38),
        size=11, bold=True, color=title_color)
    tb = slide.shapes.add_textbox(bx + Inches(0.12), by + Inches(0.57), w - Inches(0.24), h - Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"• {b}"
        r.font.size  = Pt(bullet_size)
        r.font.color.rgb = bullet_color


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  COVER
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=DARK)
rect(sl, 0, 0, Inches(0.6), H, fill=RED)
rect(sl, 0, H - Inches(2.1), W, Inches(2.1), fill=RED)
# kicker
rect(sl, Inches(0.9), Inches(0.7), Inches(3.5), Inches(0.38), fill=RED)
txt(sl, "INTERNAL PROPOSAL  ·  CONFIDENTIAL",
    Inches(0.9), Inches(0.72), Inches(3.5), Inches(0.34),
    size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# main title
txt(sl, "MUFG", Inches(0.9), Inches(1.25), Inches(9), Inches(1.0),
    size=54, bold=True, color=WHITE)
txt(sl, "Employee Benefits Portal",
    Inches(0.9), Inches(2.2), Inches(10), Inches(0.9),
    size=36, bold=True, color=WHITE)
rect(sl, Inches(0.9), Inches(3.2), Inches(4.5), Inches(0.05), fill=GOLD)
txt(sl, "A white-label B2B2C platform turning MUFG's corporate\nclient network into Indonesia's premier employee benefits channel.",
    Inches(0.9), Inches(3.4), Inches(9.5), Inches(0.9),
    size=16, color=MGREY)
# bottom meta
txt(sl, "Pilot: MUFG Indonesia employees  ·  500–1,000 headcount  ·  Q2 2026",
    Inches(0.9), H - Inches(1.85), Inches(9), Inches(0.45),
    size=13, color=WHITE, bold=True)
txt(sl, "Providers: Adira Finance  ·  Home Credit Indonesia  ·  Zurich Insurance  ·  Samsung EPP",
    Inches(0.9), H - Inches(1.38), Inches(11), Inches(0.4),
    size=12, color=WHITE)
txt(sl, "Strategy & Digital Banking Division  ·  MUFG Bank Ltd.",
    Inches(0.9), H - Inches(0.9), Inches(9), Inches(0.35),
    size=11, color=RGBColor(0xCC,0xCC,0xCC), italic=True)
# MUFG badge top right
rect(sl, W - Inches(2.1), Inches(0.2), Inches(1.8), Inches(0.8), fill=WHITE)
txt(sl, "MUFG", W - Inches(2.1), Inches(0.26), Inches(1.8), Inches(0.65),
    size=26, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, "Contents")
sections = [
    ("01", "Executive Summary"),
    ("02", "Market Opportunity"),
    ("03", "Why B2B2C — and Why Now"),
    ("04", "Solution Overview"),
    ("05", "Platform Architecture"),
    ("06", "Provider Ecosystem"),
    ("07", "Pilot Program"),
    ("08", "Go-to-Market Strategy"),
    ("09", "Product Specification"),
    ("10", "Business Case & Financials"),
    ("11", "Roadmap & Next Steps"),
]
cols = [sections[:6], sections[6:]]
for ci, col in enumerate(cols):
    for ri, (num, label) in enumerate(col):
        bx = Inches(0.5) + ci * Inches(6.5)
        by = Inches(1.25) + ri * Inches(1.0)
        rect(sl, bx, by, Inches(0.55), Inches(0.55), fill=RED)
        txt(sl, num, bx, by + Inches(0.02), Inches(0.55), Inches(0.52),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, label, bx + Inches(0.65), by + Inches(0.06), Inches(5.6), Inches(0.45),
            size=15, bold=True, color=DARK)
        rect(sl, bx + Inches(0.65), by + Inches(0.55), Inches(5.6), Inches(0.02), fill=LGREY)


# ═══════════════════════════════════════════════════════════════════════════════
# 01  EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 1, "Executive Summary",
        "One platform. Three providers. The entire MUFG corporate network.")

sl = prs.slides.add_slide(BLANK)
header(sl, "Executive Summary")

# left narrative
tb = sl.shapes.add_textbox(Inches(0.35), Inches(1.2), Inches(6.3), Inches(5.9))
tf = tb.text_frame; tf.word_wrap = True
para(tf, "What We Are Building", size=15, bold=True, color=RED, before=Pt(0))
para(tf, "A mobile-first, payroll-integrated employee benefits marketplace — white-labelled under MUFG — connecting corporate employees to vehicle financing, consumer electronics, and insurance through a single, trusted platform.", size=12, color=DGREY)
para(tf, "The Business Model", size=15, bold=True, color=RED, before=Pt(12))
para(tf, "B2B2C: MUFG sells to corporate HR clients; employees transact. Monthly instalments are deducted directly from payroll (potong gaji), eliminating repayment risk and friction simultaneously.", size=12, color=DGREY)
para(tf, "The Pilot", size=15, bold=True, color=RED, before=Pt(12))
para(tf, "500–1,000 MUFG Indonesia employees from Q2 2026, with three embedded providers from day one: Adira Finance, Home Credit Indonesia, and Zurich Insurance.", size=12, color=DGREY)
para(tf, "The Scale Play", size=15, bold=True, color=RED, before=Pt(12))
para(tf, "Extend to MUFG's full Indonesian corporate client book. Embed Samsung EPP as the anchor corporate-client marketing channel. Offer the platform as a PaaS to corporate clients by 2027.", size=12, color=DGREY)
para(tf, "The Structural Advantage", size=15, bold=True, color=RED, before=Pt(12))
para(tf, "Global precedent validates this model: Voya Financial paid $570M for Benefitfocus in 2023 precisely because bank-backed benefits distribution is a durable competitive moat. MUFG replicates this architecture natively in Indonesia.", size=12, color=DGREY)

# right stat panel
rect(sl, Inches(6.85), Inches(1.2), Inches(0.04), Inches(5.9), fill=LGREY)
stats4 = [
    ("USD 1.82B", "Global employee benefits\nplatform market (2025)"),
    ("10.4%",     "Market CAGR to 2030"),
    ("IDR 850T+", "Indonesia consumer\nfinancing market"),
    ("$570M",     "Voya/Benefitfocus deal —\nbank-backed model validated"),
]
for i, (n, l) in enumerate(stats4):
    bx = Inches(7.1)
    by = Inches(1.25) + i * Inches(1.52)
    stat_box(sl, n, l, bx, by, w=Inches(5.9), h=Inches(1.38))


# ═══════════════════════════════════════════════════════════════════════════════
# 02  MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 2, "Market Opportunity",
        "A $1.82B global market growing at 10.4% CAGR — Indonesia is the next frontier.")

sl = prs.slides.add_slide(BLANK)
header(sl, "Global Market Context")

rect(sl, Inches(0.3), Inches(1.2), W - Inches(0.6), Inches(0.78), fill=DARK)
txt(sl, "The global employee benefits platform market is USD 1.82B in 2025, projected to reach USD 2.98B by 2030 (CAGR 10.4%). North America holds ~39% of market share. Asia is the fastest-growing and most underpenetrated region.",
    Inches(0.5), Inches(1.28), W - Inches(1.0), Inches(0.62), size=12, color=WHITE)

benchmarks = [
    ("Benefitfocus\n(US)", "$570M\nacquisition", "25M+ lives served;\nacquired by Voya Financial\n2023 — bank-backed model.", RED),
    ("Benefit Systems\nMultiSport (CEE)", "EUR 2.1B\nmarket cap", "2M cards · 41,000+\nemployer partners across\n6 countries.", BLUE),
    ("Benifex\n(Europe)", "3,000+\nclients", "Formed Feb 2025 via\nBenefex + Benify merger;\n100+ countries.", DARK),
    ("Pluxee\n(Global)", "36M\nusers", "500,000+ corporate clients;\n31 countries; 3–7%\nmerchant fee model.", GREEN),
    ("Darwinbox\n(Asia)", "$72M\nraised", "Unicorn; 4M employees;\nHCM + benefits embedded\nmodel; Gartner MQ.", GOLD),
]
for i, (name, metric, detail, col) in enumerate(benchmarks):
    bx = Inches(0.3) + i * Inches(2.6)
    rect(sl, bx, Inches(2.2), Inches(2.45), Inches(4.9), fill=WHITE, line=col, lw=Pt(1.5))
    rect(sl, bx, Inches(2.2), Inches(2.45), Inches(0.5), fill=col)
    txt(sl, name, bx + Inches(0.1), Inches(2.22), Inches(2.25), Inches(0.48),
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, metric, bx + Inches(0.1), Inches(2.82), Inches(2.25), Inches(0.75),
        size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, detail, bx + Inches(0.1), Inches(3.65), Inches(2.25), Inches(1.3),
        size=10, color=DGREY, align=PP_ALIGN.CENTER)


sl = prs.slides.add_slide(BLANK)
header(sl, "Indonesia Opportunity")

opp_stats = [
    ("68M+",   "Formal-sector\nemployees — benefit\nchannel underserved"),
    ("~32%",   "Employees with\nno formal\ncredit access"),
    ("3×",     "Higher conversion\nvia employer channel\nvs. direct-to-consumer"),
    ("IDR 850T+", "Consumer financing\nmarket 2025"),
]
for i, (n, l) in enumerate(opp_stats):
    bx = Inches(0.3) + i * Inches(3.25)
    stat_box(sl, n, l, bx, Inches(1.25), w=Inches(3.0), h=Inches(2.1))

txt(sl, "Why MUFG's Payroll Channel Changes the Game",
    Inches(0.3), Inches(3.65), Inches(12.7), Inches(0.5),
    size=15, bold=True, color=RED)

drivers = [
    ("Zero repayment risk for providers",
     "Payroll deduction (potong gaji) eliminates default risk — providers can offer better rates and approve more applicants, driving higher GMV through the platform."),
    ("Pre-qualified employee base",
     "MUFG's KYC and employment data pre-qualifies employees. Providers skip costly credit bureau pulls; approval rates rise; employees experience instant decisions."),
    ("Captive distribution at near-zero CAC",
     "MUFG reaches employees via existing payroll infrastructure. Provider customer acquisition cost via this channel is estimated at 60–80% below direct consumer marketing."),
    ("Compounding network effect",
     "Each corporate client added expands both the employee user base and the co-marketing inventory available to brands like Samsung, replicating Benefit Systems' 41,000-employer flywheel in the Indonesian context."),
]
for i, (h, body) in enumerate(drivers):
    bx = Inches(0.3) + (i % 2) * Inches(6.55)
    by = Inches(4.25) + (i // 2) * Inches(1.45)
    rect(sl, bx, by, Inches(6.2), Inches(1.28), fill=LGREY)
    rect(sl, bx, by, Inches(0.07), Inches(1.28), fill=RED)
    txt(sl, h, bx + Inches(0.18), by + Inches(0.08), Inches(5.9), Inches(0.38),
        size=12, bold=True, color=RED)
    txt(sl, body, bx + Inches(0.18), by + Inches(0.48), Inches(5.9), Inches(0.72),
        size=10, color=DGREY)


# ═══════════════════════════════════════════════════════════════════════════════
# 03  WHY B2B2C — AND WHY NOW
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 3, "Why B2B2C — and Why Now",
        "Global precedent is clear. The failure patterns are equally instructive.")

sl = prs.slides.add_slide(BLANK)
header(sl, "What the Global Market Teaches Us",
       "Source: B2B2C Employee Benefits Portal Global Market Review, April 2026")

rect(sl, Inches(0.3), Inches(1.22), W - Inches(0.6), Inches(0.62), fill=DARK)
txt(sl, "Four recurring failure patterns explain every B2B2C platform collapse. MUFG's architecture is explicitly designed to avoid each one.",
    Inches(0.5), Inches(1.28), W - Inches(1.0), Inches(0.5), size=12, color=WHITE)

failures = [
    ("Failure Pattern 1\nExit-Dependent Capital",
     "Level (US, $27M raised) shut down Jan 2025 when its acquisition fell through — no path to profitability, no runway. Alza ($6.6M) followed the same month.",
     "MUFG is an institutional owner, not a venture-backed startup. The platform is a strategic business line, not an exit bet."),
    ("Failure Pattern 2\nNiche Without Cross-Sell",
     "Level served dental/vision only. Alza targeted Latino workers exclusively. Neither could expand ARPU to justify B2B sales cycle economics.",
     "MUFG's platform launches with three categories (vehicles, electronics, insurance) plus a co-marketing channel — multi-product from day one."),
    ("Failure Pattern 3\nPost-Acquisition Integration",
     "Darwin/Thomsons under Mercer: organizational chaos, talent attrition, client service degradation after the 2016 acquisition — a cautionary M&A tale.",
     "MUFG builds organically. No culture clash. Full control of product roadmap and client relationships."),
    ("Failure Pattern 4\nAdoption Debt",
     "Platforms deployed by HR and ignored by employees. Without mobile-first UX and payroll automation, enrollment rates stagnate far below contract base.",
     "Potong gaji (payroll deduction) creates inherent engagement: the platform is part of the monthly salary workflow, not an optional add-on."),
]
for i, (title, problem, mufg_ans) in enumerate(failures):
    bx = Inches(0.3) + (i % 2) * Inches(6.55)
    by = Inches(2.1) + (i // 2) * Inches(2.5)
    rect(sl, bx, by, Inches(6.2), Inches(2.25), fill=WHITE, line=RED, lw=Pt(1))
    rect(sl, bx, by, Inches(6.2), Inches(0.48), fill=RED)
    txt(sl, title, bx + Inches(0.12), by + Inches(0.07), Inches(5.96), Inches(0.38),
        size=11, bold=True, color=WHITE)
    txt(sl, f"Problem: {problem}", bx + Inches(0.12), by + Inches(0.57), Inches(5.96), Inches(0.82),
        size=10, color=DGREY, italic=True)
    rect(sl, bx + Inches(0.12), by + Inches(1.45), Inches(5.96), Inches(0.04), fill=MGREY)
    txt(sl, f"MUFG Answer: {mufg_ans}", bx + Inches(0.12), by + Inches(1.55), Inches(5.96), Inches(0.6),
        size=10, bold=True, color=DARK)


sl = prs.slides.add_slide(BLANK)
header(sl, "The Institutional-Backer Moat",
       "Voya/Benefitfocus is the closest global analogue to MUFG's strategy")

rect(sl, Inches(0.3), Inches(1.2), W - Inches(0.6), Inches(0.65), fill=DARK)
txt(sl, "\"Bank-backed benefits distribution is a durable competitive moat.\" — Global Market Review, April 2026",
    Inches(0.5), Inches(1.28), W - Inches(1.0), Inches(0.5), size=13, color=WHITE, italic=True)

comparisons = [
    ("Voya + Benefitfocus\n(US Analogue)",
     BLUE,
     ["Voya (financial institution) acquires Benefitfocus (benefits portal) for $570M",
      "Combined entity serves 38 million individuals — ~1 in 10 Americans",
      "Benefitfocus owned the employer relationship & benefits data",
      "Voya brought distribution scale and financial product depth",
      "Result: unassailable market position no pure-play startup can replicate"]),
    ("MUFG Employee Benefits Portal\n(Indonesia)",
     RED,
     ["MUFG (financial institution) builds the portal as a strategic business line — no acquisition premium",
      "MUFG owns the payroll relationship, corporate credit data, and employer trust",
      "Providers (Adira, Home Credit, Zurich, Samsung) bring products; MUFG brings distribution",
      "Potong gaji mirrors Voya's advantage: the financial institution controls the repayment rail",
      "Result: same structural moat, built from scratch at a fraction of the cost"]),
    ("Key Differentiators vs. Pure-Play Portals\n(Why MUFG Wins)",
     GREEN,
     ["No reliance on VC exit — long-term capital commitment",
      "Employer trust already established via banking relationship",
      "Payroll deduction eliminates provider repayment risk → better product rates",
      "OJK-licensed entity — no regulatory bootstrapping required",
      "Samsung EPP and future corporate client co-marketing create non-provider revenue from launch"]),
]
for i, (title, col, pts) in enumerate(comparisons):
    bx = Inches(0.3) + i * Inches(4.35)
    rect(sl, bx, Inches(2.1), Inches(4.1), Inches(5.0), fill=WHITE, line=col, lw=Pt(2))
    rect(sl, bx, Inches(2.1), Inches(4.1), Inches(0.5), fill=col)
    txt(sl, title, bx + Inches(0.12), Inches(2.14), Inches(3.86), Inches(0.44),
        size=11, bold=True, color=WHITE)
    tb = sl.shapes.add_textbox(bx + Inches(0.12), Inches(2.72), Inches(3.86), Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for pt in pts:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(5)
        r = p.add_run(); r.text = f"• {pt}"
        r.font.size = Pt(10); r.font.color.rgb = DGREY


# ═══════════════════════════════════════════════════════════════════════════════
# 04  SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 4, "Solution Overview",
        "One portal. Multiple products. Zero friction repayment.")

sl = prs.slides.add_slide(BLANK)
header(sl, "Solution: MUFG Employee Benefits Portal")

# central value proposition banner
rect(sl, Inches(0.3), Inches(1.22), W - Inches(0.6), Inches(0.72), fill=RED)
txt(sl, "A single white-label marketplace connecting MUFG's corporate employees to vehicles, devices, and insurance — repaid automatically via payroll deduction.",
    Inches(0.5), Inches(1.28), W - Inches(1.0), Inches(0.6),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

pillars = [
    ("B2B2C Ownership",
     "MUFG owns the channel. Corporate HR administers. Employees transact. Revenue flows to MUFG via origination fees, platform SaaS, and co-marketing — not from employees directly.",
     RED),
    ("Payroll Deduction\n(Potong Gaji)",
     "Zero-friction repayment deducted from monthly salary. Eliminates default risk for providers, enabling better rates and higher approval rates for employees.",
     DARK),
    ("Multi-Provider\nMarketplace",
     "Adira Finance (vehicles) · Home Credit Indonesia (devices & white goods) · Zurich Insurance (life, PA, health top-up, device protection).",
     BLUE),
    ("White-Label\nMobile App",
     "Branded for MUFG. iOS & Android. API-first architecture. Customisable per corporate client as the platform scales beyond the pilot.",
     GREEN),
    ("Corporate Marketing\nDistribution",
     "Samsung EPP embedded as a co-marketing channel from launch — exclusive employee pricing, financed via Home Credit. Template for future MUFG corporate client activations.",
     GOLD),
]
for i, (title, body, col) in enumerate(pillars):
    bx = Inches(0.3) + (i % 3) * Inches(4.35) if i < 3 else Inches(0.3) + (i-3) * Inches(6.55)
    by = Inches(2.18) if i < 3 else Inches(4.75)
    w  = Inches(4.1) if i < 3 else Inches(6.2)
    rect(sl, bx, by, w, Inches(2.35), fill=WHITE, line=col, lw=Pt(1.5))
    rect(sl, bx, by, Inches(0.06), Inches(2.35), fill=col)
    txt(sl, title, bx + Inches(0.18), by + Inches(0.1), w - Inches(0.28), Inches(0.5),
        size=12, bold=True, color=col)
    txt(sl, body, bx + Inches(0.18), by + Inches(0.65), w - Inches(0.28), Inches(1.55),
        size=11, color=DGREY)


# ═══════════════════════════════════════════════════════════════════════════════
# 05  PLATFORM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 5, "Platform Architecture",
        "API-first. Mobile-native. OJK-compliant.")

sl = prs.slides.add_slide(BLANK)
header(sl, "Platform Architecture — Layer View")

layers = [
    (RED,   "Layer 1 — Front-End",
            "White-label Mobile App (iOS / Android)  ·  Employee Web Portal  ·  HR Admin Dashboard  ·  MUFG Design System"),
    (BLUE,  "Layer 2 — Platform Core",
            "Employee Identity & SSO  ·  Product Catalogue Engine  ·  Eligibility Engine  ·  Checkout & Application Flow  ·  Notification Service"),
    (DARK,  "Layer 3 — Payroll Integration",
            "REST / SFTP to HR Payroll System  ·  Deduction Scheduler  ·  Monthly Reconciliation Engine  ·  Hardship Pause & Override"),
    (GREEN, "Layer 4 — Provider APIs",
            "Adira Finance API  ·  Home Credit Indonesia API  ·  Zurich Insurance Policy API  ·  Samsung EPP API  ·  Extensible connector framework"),
    (GOLD,  "Layer 5 — Data & Compliance",
            "OJK POJK 6/2022 (MPMM) compliance  ·  Data residency: AWS Jakarta (ap-southeast-3)  ·  Snowflake DWH  ·  Audit logs  ·  PDPA-compliant"),
]
for i, (col, label, detail) in enumerate(layers):
    by = Inches(1.25) + i * Inches(1.18)
    rect(sl, Inches(0.3), by, W - Inches(0.6), Inches(1.06), fill=WHITE, line=col, lw=Pt(1))
    rect(sl, Inches(0.3), by, Inches(2.9), Inches(1.06), fill=col)
    txt(sl, label, Inches(0.42), by + Inches(0.28), Inches(2.65), Inches(0.5),
        size=12, bold=True, color=WHITE)
    txt(sl, detail, Inches(3.35), by + Inches(0.28), Inches(9.8), Inches(0.5),
        size=11, color=DGREY)


sl = prs.slides.add_slide(BLANK)
header(sl, "Employee User Journey — End-to-End")

steps = [
    ("1\nRegister", "Download app. Authenticate via MUFG SSO. eKYC: KTP + selfie. Invite-code gated during pilot."),
    ("2\nBrowse", "Browse catalogue: Vehicles (Adira), Devices (Home Credit + Samsung EPP), Insurance (Zurich)."),
    ("3\nApply", "Select product. Auto-populated form from HR profile. Real-time eligibility check against salary."),
    ("4\nApprove", "Provider decision in <24 hours. Terms displayed in-app. E-signature for contract execution."),
    ("5\nReceive", "Vehicle delivered / device shipped / policy issued. Full digital trail retained in-app."),
    ("6\nRepay", "Automated monthly deduction from payroll. Statement in-app. Zero manual action required."),
]
aw = Inches(2.0)
for i, (title, body) in enumerate(steps):
    bx = Inches(0.28) + i * (aw + Inches(0.1))
    rect(sl, bx, Inches(1.22), aw, Inches(0.58), fill=RED)
    txt(sl, title, bx, Inches(1.24), aw, Inches(0.54),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, bx, Inches(1.8), aw, Inches(5.3), fill=LGREY, line=RED, lw=Pt(0.75))
    txt(sl, body, bx + Inches(0.1), Inches(1.95), aw - Inches(0.2), Inches(5.0),
        size=11, color=DGREY)
    if i < 5:
        txt(sl, "▶", bx + aw, Inches(1.33), Inches(0.12), Inches(0.38),
            size=13, bold=True, color=RED)

# Asia-specific insight callout
rect(sl, Inches(0.28), H - Inches(0.72), W - Inches(0.56), Inches(0.36), fill=DARK)
txt(sl, "Asia insight (Darwinbox benchmark): the winning Asian model embeds benefits in the payroll workflow — not as a standalone app. Potong gaji achieves exactly this.",
    Inches(0.4), H - Inches(0.70), W - Inches(0.8), Inches(0.3),
    size=9, color=WHITE, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 06  PROVIDER ECOSYSTEM
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 6, "Provider Ecosystem",
        "Three providers at launch. One corporate marketing channel.")

sl = prs.slides.add_slide(BLANK)
header(sl, "Day-1 Provider Ecosystem")

providers = [
    ("ADIRA FINANCE", RED,
     ["Category: Vehicle Financing",
      "Products: New & used car loans; motorcycle loans",
      "Tenor: 12–60 months  ·  DP: from 10%",
      "Pricing: fleet/group rate (below retail)",
      "Integration: REST API · eKYC · instant scoring",
      "Revenue to MUFG: origination fee 0.5–1.0% per deal"]),
    ("HOME CREDIT INDONESIA", BLUE,
     ["Category: Consumer Electronics & White Goods",
      "Products: Smartphones (incl. Samsung EPP), laptops, TVs, appliances",
      "Tenor: 3–36 months  ·  0% instalment promos",
      "No collateral; instant POS approval",
      "Integration: POS API · merchant catalogue sync",
      "Revenue to MUFG: origination fee + Samsung co-marketing"]),
    ("ZURICH INSURANCE", DARK,
     ["Category: Insurance",
      "Products: Term life (group), PA, health top-up, device protection",
      "Premium via payroll deduction — group rates 20–30% below retail",
      "Digital claims via app",
      "Integration: REST API · policy issuance · claims API",
      "Revenue to MUFG: origination fee + referral commission"]),
]
for i, (name, col, pts) in enumerate(providers):
    bx = Inches(0.3) + i * Inches(4.35)
    rect(sl, bx, Inches(1.22), Inches(4.1), Inches(5.9), fill=WHITE, line=col, lw=Pt(2))
    rect(sl, bx, Inches(1.22), Inches(4.1), Inches(0.55), fill=col)
    txt(sl, name, bx + Inches(0.12), Inches(1.26), Inches(3.86), Inches(0.48),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = sl.shapes.add_textbox(bx + Inches(0.15), Inches(1.92), Inches(3.8), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for pt in pts:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(5)
        r = p.add_run(); r.text = f"• {pt}"
        r.font.size = Pt(11); r.font.color.rgb = DGREY


sl = prs.slides.add_slide(BLANK)
header(sl, "Corporate Marketing Channel: Samsung EPP",
       "The platform as distribution — not just benefits administration")

rect(sl, Inches(0.3), Inches(1.22), W - Inches(0.6), Inches(0.65), fill=DARK)
txt(sl, "Beyond provider fees: MUFG's portal is a co-marketing platform for corporate clients' consumer products. Samsung EPP is the Day-1 anchor.",
    Inches(0.5), Inches(1.28), W - Inches(1.0), Inches(0.52),
    size=12, color=WHITE, align=PP_ALIGN.CENTER)

epp_cols = [
    ("What is Samsung EPP?",
     ["Employee Purchase Programme — exclusive Samsung discounts for corporate employees",
      "Galaxy smartphones, tablets, wearables at 10–25% below MSRP",
      "Financed via 0% instalment through Home Credit Indonesia (already a platform partner — no extra integration)",
      "Employee buys in-app; instalment deducted via payroll; Samsung ships direct",
      "Samsung Indonesia gains a low-CAC B2B2C channel with pre-qualified, salaried consumers"]),
    ("Revenue for MUFG",
     ["Co-marketing fee from Samsung Indonesia per campaign",
      "Transaction fee per order processed through the portal",
      "Increased platform stickiness: Samsung activations drive monthly app opens",
      "Proof-of-concept for other MUFG corporate clients (automotive OEMs, FMCG brands, telcos)",
      "Expected revenue: IDR 200–500M per Samsung campaign cycle"]),
    ("Why It Works — The Benefit Systems Analogy",
     ["Benefit Systems (Poland) built EUR 2.1B market cap on a single product: MultiSport gym access",
      "Network effect: 41,000+ employers → 2M users → gyms compete to join the network",
      "MUFG replicates this logic: Samsung EPP is the anchor; other brands follow to access the employee base",
      "Each corporate client added increases the co-marketing inventory → rising platform value for all parties",
      "This is the 'second revenue engine' that pure-play benefits portals lack"]),
]
for i, (title, pts) in enumerate(epp_cols):
    bx = Inches(0.3) + i * Inches(4.35)
    rect(sl, bx, Inches(2.1), Inches(4.1), Inches(5.0), fill=WHITE, line=RED, lw=Pt(1))
    rect(sl, bx, Inches(2.1), Inches(4.1), Inches(0.48), fill=RED)
    txt(sl, title, bx + Inches(0.12), Inches(2.14), Inches(3.86), Inches(0.4),
        size=11, bold=True, color=WHITE)
    tb = sl.shapes.add_textbox(bx + Inches(0.12), Inches(2.68), Inches(3.86), Inches(4.3))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for pt in pts:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(5)
        r = p.add_run(); r.text = f"• {pt}"
        r.font.size = Pt(10); r.font.color.rgb = DGREY


# ═══════════════════════════════════════════════════════════════════════════════
# 07  PILOT PROGRAM
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 7, "Pilot Program",
        "MUFG Indonesia  ·  500–1,000 employees  ·  Q2–Q3 2026")

sl = prs.slides.add_slide(BLANK)
header(sl, "Pilot Program: Design & Success Criteria")

pilot_cards = [
    ("Participants", RED,
     ["500–1,000 MUFG Indonesia permanent staff",
      "All employee grades included",
      "HR lead: MUFG Indonesia HR Division",
      "Pilot window: Q2–Q3 2026 (6 months)"]),
    ("Products Offered", BLUE,
     ["Adira: Motorcycle (Honda/Yamaha) + car loans",
      "Home Credit: Samsung Galaxy bundles + white goods",
      "Zurich: Group term life, PA plan, device protection",
      "Samsung EPP: Exclusive pricing soft-launch"]),
    ("Success Metrics", DARK,
     ["≥ 30% employee registration within 90 days",
      "≥ 15% transaction conversion (≥1 product)",
      "Zero payroll reconciliation errors",
      "Employee NPS ≥ 40"]),
    ("Operational Setup", GREEN,
     ["Payroll integration: API with MUFG HR payroll system",
      "Onboarding: eKYC (KTP + selfie + liveness check)",
      "Deduction schedule: monthly, aligned to payroll run",
      "Support: in-app chat + dedicated HR helpdesk"]),
]
for i, (title, col, pts) in enumerate(pilot_cards):
    bx = Inches(0.3) + (i % 2) * Inches(6.55)
    by = Inches(1.25) + (i // 2) * Inches(2.95)
    card(sl, title, pts, bx, by, Inches(6.2), Inches(2.7),
         title_bg=col, bullet_size=12)

# Asia benchmark note
rect(sl, Inches(0.3), H - Inches(0.72), W - Inches(0.56), Inches(0.36), fill=LGREY)
txt(sl, "Asia context: 49% of APAC organisations cite meeting employee expectations as their #2 challenge (up 11pp YoY — CIPD 2024). The pilot's NPS target directly benchmarks against this gap.",
    Inches(0.45), H - Inches(0.70), W - Inches(0.8), Inches(0.3),
    size=9, color=DGREY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 08  GO-TO-MARKET
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 8, "Go-to-Market Strategy",
        "Pilot → Scale → Platform-as-a-Service")

sl = prs.slides.add_slide(BLANK)
header(sl, "Go-to-Market: Three Phases")

phases = [
    ("Phase 1\nQ1–Q2 2026\nBuild & Integrate", RED,
     ["Build platform MVP (mobile app + web portal + APIs)",
      "Onboard 3 providers (Adira, Home Credit, Zurich)",
      "Sign Samsung EPP co-marketing MOU",
      "HR payroll API integration with MUFG Indonesia",
      "OJK compliance filing & security audit"]),
    ("Phase 2\nQ3–Q4 2026\nPilot & Learn", BLUE,
     ["Launch pilot: 500–1,000 MUFG Indonesia employees",
      "Samsung EPP soft-launch on platform",
      "Collect usage data, iterate UX weekly",
      "HR admin dashboard v1 live",
      "Prepare board report on pilot outcomes"]),
    ("Phase 3\nQ1–Q2 2027\nScale & PaaS", DARK,
     ["Expand to 5 MUFG corporate clients (target: 10,000 employees)",
      "Add 2–3 new providers (BNPL, travel insurance)",
      "Samsung EPP full national rollout",
      "Launch Platform-as-a-Service for MUFG corporate clients",
      "Target: 25,000+ employees on platform by end of 2027"]),
]
for i, (phase, col, pts) in enumerate(phases):
    bx = Inches(0.3) + i * Inches(4.35)
    rect(sl, bx, Inches(1.22), Inches(4.1), Inches(5.9), fill=col)
    txt(sl, phase, bx + Inches(0.15), Inches(1.32), Inches(3.8), Inches(1.1),
        size=14, bold=True, color=WHITE)
    tb = sl.shapes.add_textbox(bx + Inches(0.15), Inches(2.55), Inches(3.8), Inches(4.4))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for pt in pts:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(7)
        r = p.add_run(); r.text = f"• {pt}"
        r.font.size = Pt(12); r.font.color.rgb = WHITE


# ═══════════════════════════════════════════════════════════════════════════════
# 09  PRODUCT SPECIFICATION
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 9, "Product Specification",
        "Mobile-first. API-first. Compliance-built-in.")

sl = prs.slides.add_slide(BLANK)
header(sl, "Product Specification — App Feature Set")

features = [
    ("Authentication & Onboarding",
     ["SSO via MUFG employee directory",
      "eKYC: KTP + selfie + liveness check",
      "Biometric login (Face ID / Fingerprint)",
      "Invite-code gated during pilot"], RED),
    ("Product Catalogue",
     ["Category tiles: Vehicles · Devices · Insurance · Samsung EPP",
      "Real-time inventory & availability",
      "Personalised eligibility filter",
      "Exclusive employee pricing badges"], BLUE),
    ("Application Flow",
     ["In-app form auto-populated from HR profile",
      "Real-time eligibility engine (salary vs instalment ratio)",
      "Document upload: NPWP, KTP, salary slip",
      "E-signature for contract execution"], DARK),
    ("Payroll Deduction",
     ["Deduction schedule synced with monthly payroll",
      "Employee consent & deduction dashboard",
      "HR override & hardship pause capability",
      "Monthly statement & repayment history"], GREEN),
    ("Notifications & Support",
     ["Push: approval, deduction, promotions",
      "In-app chat (CS) + HR helpdesk escalation",
      "Zurich insurance claims submission",
      "Referral & loyalty points (Phase 2)"], GOLD),
    ("Admin Dashboard",
     ["HR: headcount, deduction status, reconciliation",
      "Provider: leads, approvals, disbursements",
      "MUFG: revenue, GMV, platform health KPIs",
      "OJK-compliant audit trail & reporting"], RED),
]
for i, (title, pts, col) in enumerate(features):
    bx = Inches(0.3) + (i % 3) * Inches(4.35)
    by = Inches(1.25) + (i // 3) * Inches(3.05)
    card(sl, title, pts, bx, by, Inches(4.1), Inches(2.8),
         title_bg=col, bullet_size=11)


sl = prs.slides.add_slide(BLANK)
header(sl, "Technical Architecture")

tech = [
    (RED,   "Front-End",           "React Native (iOS + Android)  ·  Next.js Web Portal  ·  MUFG Design System"),
    (BLUE,  "API Gateway",         "Kong / AWS API Gateway  ·  OAuth 2.0 / OpenID Connect  ·  WAF + rate limiting"),
    (DARK,  "Back-End",            "Node.js microservices  ·  Eligibility Engine (Python ML)  ·  Notification Service (Firebase)"),
    (GREEN, "Payroll Connector",   "REST / SFTP to HR system  ·  Monthly deduction scheduler  ·  Reconciliation worker"),
    (GOLD,  "Provider Integrations","Adira REST API  ·  Home Credit REST API  ·  Zurich Policy API  ·  Samsung EPP API"),
    (RED,   "Data & Analytics",    "PostgreSQL + Redis  ·  Snowflake DWH  ·  Metabase dashboards  ·  Segment CDP"),
    (BLUE,  "Infrastructure",      "AWS Jakarta (ap-southeast-3)  ·  Kubernetes (EKS)  ·  Terraform IaC  ·  SOC 2 Type II"),
    (DARK,  "Compliance",          "OJK POJK 6/2022 (MPMM)  ·  Data residency: Indonesia  ·  PDPA  ·  ISO 27001 roadmap"),
]
for i, (col, label, detail) in enumerate(tech):
    by = Inches(1.28) + i * Inches(0.74)
    rect(sl, Inches(0.3), by, Inches(2.6), Inches(0.65), fill=col)
    txt(sl, label, Inches(0.38), by + Inches(0.1), Inches(2.42), Inches(0.46),
        size=11, bold=True, color=WHITE)
    rect(sl, Inches(2.9), by, Inches(10.1), Inches(0.65), fill=LGREY if i % 2 == 0 else WHITE, line=col, lw=Pt(0.5))
    txt(sl, detail, Inches(3.0), by + Inches(0.1), Inches(9.9), Inches(0.46),
        size=10.5, color=DGREY)


# ═══════════════════════════════════════════════════════════════════════════════
# 10  BUSINESS CASE & FINANCIALS
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 10, "Business Case & Financials",
        "Payback ~18 months  ·  EBITDA-positive Year 2  ·  IRR > 40% over 5 years")

sl = prs.slides.add_slide(BLANK)
header(sl, "Revenue Model")

streams = [
    ("Origination Fee",      "0.5–1.5% of financed amount, charged to provider per approved application.",
     "Primary recurring revenue. Scales linearly with GMV.", RED),
    ("Platform SaaS Fee",    "IDR 5–15M per corporate client per month.",
     "Scales with number of employer clients onboarded.", BLUE),
    ("Co-marketing Revenue", "IDR 50–200M per campaign (Samsung EPP anchor).",
     "Non-linear upside; each new corporate client adds inventory.", DARK),
    ("Interest Income",      "3–5% p.a. margin on MUFG-funded instalment products (Phase 2+).",
     "Future balance-sheet revenue — not in Year 1 projections.", GREEN),
    ("Data & Analytics",     "IDR 10–50M per anonymised workforce analytics report.",
     "High-margin add-on; sold to providers and HR consultancies.", GOLD),
]
for i, (name, desc, strategic, col) in enumerate(streams):
    bx = Inches(0.3) + (i % 3) * Inches(4.35) if i < 3 else Inches(0.3) + (i-3) * Inches(6.55)
    by = Inches(1.25) if i < 3 else Inches(4.35)
    w  = Inches(4.1) if i < 3 else Inches(6.2)
    rect(sl, bx, by, w, Inches(2.8), fill=WHITE, line=col, lw=Pt(1.5))
    rect(sl, bx, by, Inches(0.06), Inches(2.8), fill=col)
    txt(sl, name,     bx + Inches(0.18), by + Inches(0.12), w - Inches(0.28), Inches(0.42),
        size=13, bold=True, color=col)
    txt(sl, desc,     bx + Inches(0.18), by + Inches(0.58), w - Inches(0.28), Inches(0.9),
        size=11, color=DGREY)
    rect(sl, bx + Inches(0.18), by + Inches(1.55), w - Inches(0.28), Inches(0.04), fill=MGREY)
    txt(sl, strategic, bx + Inches(0.18), by + Inches(1.65), w - Inches(0.28), Inches(0.95),
        size=10, color=DGREY, italic=True)


sl = prs.slides.add_slide(BLANK)
header(sl, "3-Year Financial Projection  (IDR Billion)")

col_x = [Inches(0.3), Inches(4.3), Inches(6.95), Inches(9.6)]
col_w = [Inches(3.9), Inches(2.55), Inches(2.55), Inches(2.55)]
hdrs  = ["  Item", "Year 1 — Pilot\n(2026)", "Year 2 — Scale\n(2027)", "Year 3 — Growth\n(2028)"]
for j in range(4):
    rect(sl, col_x[j], Inches(1.25), col_w[j], Inches(0.7), fill=DARK)
    txt(sl, hdrs[j], col_x[j] + Inches(0.08), Inches(1.28), col_w[j] - Inches(0.16), Inches(0.62),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

rows = [
    ("Employees on Platform",             "750",      "8,000",    "25,000"),
    ("Corporate Clients",                  "1",        "5",        "15"),
    ("GMV  (IDR Bn)",                      "7.5",      "80",       "250"),
    ("Origination Fee Revenue",            "0.90",     "9.60",     "30.00"),
    ("Platform SaaS Revenue",              "0.06",     "0.36",     "1.08"),
    ("Co-marketing Revenue",               "0.20",     "0.60",     "1.50"),
    ("Total Revenue",                      "1.16",     "10.56",    "32.58"),
    ("Platform OPEX",                      "(1.80)",   "(3.50)",   "(6.00)"),
    ("CAPEX  (amortised)",                 "(1.20)",   "(0.60)",   "(0.40)"),
    ("EBITDA",                             "(1.84)",   "6.46",     "26.18"),
    ("EBITDA Margin",                      "—",        "61%",      "80%"),
]
row_bg  = [LGREY,WHITE,LGREY,WHITE,LGREY,WHITE,
           RGBColor(0xFF,0xE8,0xE8), LGREY, WHITE,
           RED, LGREY]
row_fg  = [DGREY]*9 + [WHITE, DGREY]

for ri, (label, y1, y2, y3) in enumerate(rows):
    by = Inches(1.95) + ri * Inches(0.47)
    vals = [f"  {label}", y1, y2, y3]
    bold = ri in (6, 9)
    for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
        rect(sl, cx, by, cw, Inches(0.44), fill=row_bg[ri])
        txt(sl, val, cx + Inches(0.05), by + Inches(0.05), cw - Inches(0.1), Inches(0.36),
            size=10, bold=bold, color=row_fg[ri],
            align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

txt(sl, "* Assumptions: 15% transaction conversion rate · Avg ticket IDR 10M · 3-provider origination fee model · Samsung EPP co-marketing included Year 1",
    Inches(0.3), H - Inches(0.62), W - Inches(0.6), Inches(0.25),
    size=8, color=MGREY, italic=True)


sl = prs.slides.add_slide(BLANK)
header(sl, "Investment Required & Return Profile")

inv = [
    ("Platform Build  (MVP)", "IDR 8–10 Billion",
     "Mobile app, web portal, API integrations, payroll connector, security audit."),
    ("Pilot Operations  (6 months)", "IDR 1.5 Billion",
     "3 FTE product & ops staff, cloud infrastructure, provider onboarding, pilot marketing."),
    ("Compliance & Licensing", "IDR 0.5 Billion",
     "OJK registration, legal structuring, data-protection audits, PDPA compliance."),
]
for i, (label, amount, detail) in enumerate(inv):
    by = Inches(1.3) + i * Inches(1.2)
    rect(sl, Inches(0.3), by, W - Inches(0.6), Inches(1.05), fill=LGREY, line=RED, lw=Pt(1))
    rect(sl, Inches(0.3), by, Inches(0.06), Inches(1.05), fill=RED)
    txt(sl, label,  Inches(0.5),  by + Inches(0.1),  Inches(7.5), Inches(0.42),
        size=13, bold=True, color=DARK)
    txt(sl, amount, Inches(10.3), by + Inches(0.06), Inches(2.8), Inches(0.52),
        size=20, bold=True, color=RED, align=PP_ALIGN.RIGHT)
    txt(sl, detail, Inches(0.5),  by + Inches(0.57), Inches(9.5), Inches(0.38),
        size=11, color=DGREY, italic=True)

# Total
rect(sl, Inches(0.3), Inches(4.95), W - Inches(0.6), Inches(0.88), fill=RED)
txt(sl, "TOTAL INVESTMENT", Inches(0.5), Inches(5.0), Inches(7.5), Inches(0.55),
    size=16, bold=True, color=WHITE)
txt(sl, "IDR 10–12 Billion", Inches(8.0), Inches(4.97), Inches(5.0), Inches(0.68),
    size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# Return panel
rect(sl, Inches(0.3), Inches(6.0), W - Inches(0.6), Inches(1.1), fill=DARK)
returns = [
    ("Payback Period", "~18 months"),
    ("Break-Even", "End of Year 2"),
    ("5-Year IRR", "> 40%"),
    ("5-Year NPV", "IDR 85–110 Billion"),
]
for i, (label, val) in enumerate(returns):
    bx = Inches(0.5) + i * Inches(3.2)
    rect(sl, bx, Inches(6.08), Inches(3.0), Inches(0.88), fill=DARK)
    txt(sl, label, bx, Inches(6.1), Inches(3.0), Inches(0.35),
        size=10, color=MGREY, align=PP_ALIGN.CENTER)
    txt(sl, val,   bx, Inches(6.44), Inches(3.0), Inches(0.48),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 11  ROADMAP & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
divider(prs, 11, "Roadmap & Next Steps",
        "From pilot approval to PaaS platform — a 5-quarter path.")

sl = prs.slides.add_slide(BLANK)
header(sl, "Implementation Roadmap")

milestones = [
    ("Q1 2026",
     ["Platform design & architecture", "Provider contract negotiations",
      "HR payroll API discovery", "Samsung EPP MOU signing"]),
    ("Q2 2026",
     ["MVP development (mobile + API)", "Payroll integration build",
      "eKYC integration (3 providers)", "Internal UAT & security audit"]),
    ("Q3 2026",
     ["PILOT LAUNCH — 500–1,000 MUFG\nIndonesia employees",
      "Samsung EPP soft-launch", "Weekly ops review cadence",
      "UX iteration from usage data"]),
    ("Q4 2026",
     ["Pilot results & board report",
      "Corporate clients 2–5 onboarding",
      "OJK reporting framework live", "Phase 2 planning & budgeting"]),
    ("Q1 2027",
     ["Scale: 5 clients · 10,000 employees",
      "BNPL & travel insurance providers",
      "Samsung EPP national rollout",
      "PaaS offering for MUFG clients"]),
]
for i, (qtr, tasks) in enumerate(milestones):
    bx = Inches(0.3) + i * Inches(2.6)
    is_launch = i == 2
    rect(sl, bx, Inches(1.22), Inches(2.45), Inches(0.58),
         fill=GOLD if is_launch else RED)
    txt(sl, qtr, bx, Inches(1.25), Inches(2.45), Inches(0.52),
        size=13, bold=True, color=DARK if is_launch else WHITE, align=PP_ALIGN.CENTER)
    rect(sl, bx, Inches(1.8), Inches(2.45), Inches(5.3),
         fill=RGBColor(0xFF,0xF8,0xE8) if is_launch else LGREY, line=GOLD if is_launch else RED, lw=Pt(1))
    tb = sl.shapes.add_textbox(bx + Inches(0.1), Inches(1.9), Inches(2.28), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for t in tasks:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(6)
        r = p.add_run(); r.text = f"• {t}"
        r.font.size = Pt(10)
        r.font.color.rgb = DARK if is_launch else DGREY
        r.font.bold = is_launch and tasks.index(t) == 0


sl = prs.slides.add_slide(BLANK)
header(sl, "Immediate Next Steps — Action Plan")

steps = [
    ("01", "Executive Alignment\nJanuary 2026",
     ["Present proposal to MUFG Indonesia CEO & COO",
      "Align Digital Banking and Retail divisions on ownership",
      "Confirm FY2026 budget envelope",
      "Appoint internal Product Owner"]),
    ("02", "Provider Letters of Intent\nJanuary–February 2026",
     ["Issue LOI to Adira Finance",
      "Issue LOI to Home Credit Indonesia",
      "Issue LOI to Zurich Insurance Indonesia",
      "Begin Samsung EPP commercial term negotiation"]),
    ("03", "Technical Discovery\nFebruary 2026",
     ["HR payroll API specification review",
      "Provider API documentation deep-dive",
      "OJK POJK 6/2022 compliance gap assessment",
      "Vendor RFP for platform development partner"]),
    ("04", "Board Approval\nMarch 2026",
     ["Full business case for board sign-off",
      "Budget approval: IDR 10–12 Billion capex",
      "Legal entity structure for platform ownership",
      "Formal pilot agreement with MUFG Indonesia HR"]),
]
for i, (num, title, tasks) in enumerate(steps):
    bx = Inches(0.3) + (i % 2) * Inches(6.55)
    by = Inches(1.25) + (i // 2) * Inches(2.95)
    rect(sl, bx, by, Inches(6.2), Inches(2.7), fill=WHITE, line=RED, lw=Pt(1.5))
    rect(sl, bx, by, Inches(0.72), Inches(2.7), fill=RED)
    txt(sl, num, bx + Inches(0.01), by + Inches(0.75), Inches(0.7), Inches(1.0),
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, bx + Inches(0.82), by + Inches(0.1), Inches(5.2), Inches(0.6),
        size=12, bold=True, color=RED)
    tb = sl.shapes.add_textbox(bx + Inches(0.82), by + Inches(0.78), Inches(5.2), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for t in tasks:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(4)
        r = p.add_run(); r.text = f"• {t}"
        r.font.size = Pt(11); r.font.color.rgb = DGREY


# ═══════════════════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=DARK)
rect(sl, 0, 0, Inches(0.6), H, fill=RED)
rect(sl, 0, H - Inches(2.0), W, Inches(2.0), fill=RED)

txt(sl, "The Platform That Turns", Inches(0.9), Inches(1.3), Inches(11), Inches(0.65),
    size=26, color=MGREY)
txt(sl, "Payroll into Benefits.", Inches(0.9), Inches(1.9), Inches(11), Inches(1.1),
    size=46, bold=True, color=WHITE)
rect(sl, Inches(0.9), Inches(3.12), Inches(4.5), Inches(0.06), fill=GOLD)
txt(sl, "Adira Finance  ·  Home Credit Indonesia  ·  Zurich Insurance  ·  Samsung EPP",
    Inches(0.9), Inches(3.28), Inches(11), Inches(0.48),
    size=13, color=MGREY)
txt(sl, "Powered by MUFG Bank Ltd.",
    Inches(0.9), Inches(3.85), Inches(6), Inches(0.42),
    size=14, bold=True, color=WHITE)

txt(sl, "CONFIDENTIAL  ·  For internal discussion only  ·  © 2026 MUFG Bank Ltd.",
    Inches(0.9), H - Inches(1.75), Inches(11), Inches(0.4),
    size=12, color=WHITE, bold=True)
txt(sl, "Strategy & Digital Banking Division  ·  MUFG Indonesia",
    Inches(0.9), H - Inches(1.28), Inches(11), Inches(0.38),
    size=11, color=WHITE)
txt(sl, "Market data sourced from: B2B2C Employee Benefits Portal Global Market Review, April 2026",
    Inches(0.9), H - Inches(0.82), Inches(11), Inches(0.32),
    size=9, color=RGBColor(0xCC,0xCC,0xCC), italic=True)

# MUFG badge
rect(sl, W - Inches(2.1), Inches(0.22), Inches(1.8), Inches(0.78), fill=WHITE)
txt(sl, "MUFG", W - Inches(2.1), Inches(0.28), Inches(1.8), Inches(0.65),
    size=26, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "/Users/msugirin/ClaudeWS/B2B2BC-Project/MUFG_Employee_Benefits_Portal_v2.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
