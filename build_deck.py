"""
MUFG Employee Benefits Portal — Bilingual (EN/BI) Pitch Deck Generator
Produces a single PPTX covering: Pitch Deck + Product Spec + Business Case
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Palette ─────────────────────────────────────────────────────────────
MUFG_RED    = RGBColor(0xC8, 0x10, 0x2E)   # MUFG signature red
MUFG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
MUFG_GOLD   = RGBColor(0xB8, 0x86, 0x0B)   # dark goldenrod accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
MID_GREY    = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GREY   = RGBColor(0x44, 0x44, 0x44)
ACCENT_BLUE = RGBColor(0x00, 0x5F, 0xAF)   # secondary accent

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper primitives ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=14, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, space_before=Pt(6), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def bullet_box(slide, lines, x, y, w, h,
               font_size=13, color=DARK_GREY, title=None, title_color=MUFG_RED):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    if title:
        add_para(tf, title, font_size=font_size+1, bold=True,
                 color=title_color, space_before=Pt(0))
    first = True
    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2) if not first else Pt(4)
        first = False
        run = p.add_run()
        run.text = ("• " if not line.startswith("–") else "") + line
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
    return txb


# ── Section divider helper ────────────────────────────────────────────────────

def section_divider(prs, section_en, section_bi, num):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=MUFG_DARK)
    add_rect(sl, 0, 0, Inches(0.5), H, fill=MUFG_RED)
    add_text(sl, f"SECTION {num}", Inches(1), Inches(2.2), Inches(11), Inches(0.6),
             font_size=14, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)
    add_text(sl, section_en, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
             font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(sl, section_bi, Inches(1), Inches(3.9), Inches(11), Inches(0.8),
             font_size=22, bold=False, color=MUFG_GOLD, align=PP_ALIGN.LEFT, italic=True)
    # thin gold rule
    add_rect(sl, Inches(1), Inches(3.75), Inches(4), Inches(0.03), fill=MUFG_GOLD)
    return sl


# ── Slide header bar ──────────────────────────────────────────────────────────

def slide_header(slide, title_en, subtitle_bi="", bg=LIGHT_GREY):
    add_rect(slide, 0, 0, W, H, fill=bg)
    add_rect(slide, 0, 0, W, Inches(1.1), fill=MUFG_RED)
    add_rect(slide, 0, Inches(1.1), Inches(0.07), H - Inches(1.1), fill=MUFG_RED)
    add_text(slide, title_en, Inches(0.3), Inches(0.1), Inches(9.5), Inches(0.75),
             font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_bi:
        add_text(slide, subtitle_bi, Inches(0.3), Inches(0.75), Inches(9.5), Inches(0.4),
                 font_size=13, bold=False, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.LEFT, italic=True)
    # MUFG badge top-right
    add_rect(slide, W - Inches(1.8), Inches(0.15), Inches(1.6), Inches(0.75), fill=WHITE)
    add_text(slide, "MUFG", W - Inches(1.8), Inches(0.18), Inches(1.6), Inches(0.5),
             font_size=20, bold=True, color=MUFG_RED, align=PP_ALIGN.CENTER)
    # footer bar
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=MUFG_DARK)
    add_text(slide, "CONFIDENTIAL  |  MUFG Employee Benefits Portal  |  Internal Proposal 2026",
             Inches(0.3), H - Inches(0.32), W - Inches(0.6), Inches(0.28),
             font_size=8, color=MID_GREY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=MUFG_DARK)
add_rect(sl, 0, 0, Inches(0.7), H, fill=MUFG_RED)
add_rect(sl, 0, H - Inches(2.2), W, Inches(2.2), fill=MUFG_RED)

add_text(sl, "MUFG", Inches(1.1), Inches(0.6), Inches(5), Inches(0.9),
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Employee Benefits Portal", Inches(1.1), Inches(1.45), Inches(11), Inches(0.9),
         font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Portal Manfaat Karyawan MUFG", Inches(1.1), Inches(2.25), Inches(11), Inches(0.6),
         font_size=20, bold=False, color=MUFG_GOLD, align=PP_ALIGN.LEFT, italic=True)

add_rect(sl, Inches(1.1), Inches(2.95), Inches(5), Inches(0.04), fill=MUFG_GOLD)

add_text(sl, "White-Label B2B2C Benefits & Financing Platform", Inches(1.1), Inches(3.1), Inches(10), Inches(0.5),
         font_size=16, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)
add_text(sl, "Platform Manfaat & Pembiayaan B2B2C White-Label", Inches(1.1), Inches(3.55), Inches(10), Inches(0.5),
         font_size=14, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT, italic=True)

# bottom panel
add_text(sl, "INTERNAL PROPOSAL  |  CONFIDENTIAL", Inches(1.1), H - Inches(2.0), Inches(7), Inches(0.4),
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Pilot: MUFG Indonesia  |  Q2 2026", Inches(1.1), H - Inches(1.6), Inches(7), Inches(0.4),
         font_size=13, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Prepared by: Strategy & Digital Banking", Inches(1.1), H - Inches(1.2), Inches(7), Inches(0.4),
         font_size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# provider logos text placeholders
for i, (name, col) in enumerate([("Adira Finance", MUFG_GOLD),
                                   ("Home Credit Indonesia", WHITE),
                                   ("Zurich Insurance", WHITE)]):
    add_rect(sl, W - Inches(3.6) + i * Inches(1.22) - Inches(1.22*2), H - Inches(1.9),
             Inches(1.1), Inches(0.5), line=WHITE)
    add_text(sl, name, W - Inches(3.6) + i * Inches(1.22) - Inches(1.22*2), H - Inches(1.9),
             Inches(1.1), Inches(0.5), font_size=7, bold=True, color=col, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Table of Contents  |  Daftar Isi")
sections = [
    ("01", "Executive Summary", "Ringkasan Eksekutif"),
    ("02", "Problem & Opportunity", "Masalah & Peluang"),
    ("03", "Solution Overview", "Gambaran Solusi"),
    ("04", "Platform Architecture", "Arsitektur Platform"),
    ("05", "Provider Ecosystem", "Ekosistem Mitra"),
    ("06", "Pilot Program", "Program Percontohan"),
    ("07", "Go-to-Market & Distribution", "Strategi Distribusi"),
    ("08", "Product Specification", "Spesifikasi Produk"),
    ("09", "Business Case & Financials", "Kasus Bisnis & Keuangan"),
    ("10", "Roadmap & Next Steps", "Peta Jalan & Langkah Selanjutnya"),
]
cols = [sections[:5], sections[5:]]
for ci, col in enumerate(cols):
    for ri, (num, en, bi) in enumerate(col):
        bx = Inches(0.5) + ci * Inches(6.6)
        by = Inches(1.4) + ri * Inches(1.1)
        add_rect(sl, bx, by, Inches(0.5), Inches(0.5), fill=MUFG_RED)
        add_text(sl, num, bx, by, Inches(0.5), Inches(0.5),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, en, bx + Inches(0.55), by, Inches(5.8), Inches(0.32),
                 font_size=14, bold=True, color=MUFG_DARK, align=PP_ALIGN.LEFT)
        add_text(sl, bi, bx + Inches(0.55), by + Inches(0.3), Inches(5.8), Inches(0.28),
                 font_size=11, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Executive Summary", "Ringkasan Eksekutif", "01")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Executive Summary  |  Ringkasan Eksekutif")
# left column — EN
txb = sl.shapes.add_textbox(Inches(0.35), Inches(1.3), Inches(6.1), Inches(5.7))
tf = txb.text_frame; tf.word_wrap = True
add_para(tf, "The Opportunity", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(0))
add_para(tf, "MUFG Indonesia sits at the intersection of two high-growth markets: employee financial wellness and consumer financing. By launching a white-label Employee Benefits Portal, MUFG transforms its corporate client base into a powerful B2B2C distribution channel.", font_size=12, color=DARK_GREY)
add_para(tf, "What We Are Building", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(10))
add_para(tf, "A mobile-first, payroll-integrated marketplace connecting MUFG-employed and MUFG-served corporate employees to curated financial products — vehicle financing (Adira), consumer electronics & white goods (Home Credit Indonesia), and insurance (Zurich).", font_size=12, color=DARK_GREY)
add_para(tf, "The Pilot", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(10))
add_para(tf, "500–1,000 MUFG Indonesia employees  |  Q2 2026  |  Repayment via potong gaji (payroll deduction).", font_size=12, color=DARK_GREY)
add_para(tf, "Scale-Up Vision", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(10))
add_para(tf, "Extend the platform to MUFG's entire Indonesian corporate client book and introduce Samsung Electronics Partnership Program (EPP) as the anchor corporate-client distribution example.", font_size=12, color=DARK_GREY)

# right column — BI
txb2 = sl.shapes.add_textbox(Inches(6.7), Inches(1.3), Inches(6.1), Inches(5.7))
tf2 = txb2.text_frame; tf2.word_wrap = True
add_para(tf2, "Peluang", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(0))
add_para(tf2, "MUFG Indonesia berada di persimpangan dua pasar berdaya tumbuh tinggi: kesejahteraan keuangan karyawan dan pembiayaan konsumen. Dengan meluncurkan Portal Manfaat Karyawan bermerek sendiri (white-label), MUFG mengubah basis klien korporatnya menjadi saluran distribusi B2B2C yang andal.", font_size=12, color=DARK_GREY, italic=True)
add_para(tf2, "Yang Kami Bangun", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(10))
add_para(tf2, "Marketplace berbasis mobile, terintegrasi dengan payroll, yang menghubungkan karyawan dengan produk keuangan terkurasi — pembiayaan kendaraan (Adira), elektronik & barang kebutuhan rumah tangga (Home Credit), dan asuransi (Zurich).", font_size=12, color=DARK_GREY, italic=True)
add_para(tf2, "Program Percontohan", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(10))
add_para(tf2, "500–1.000 karyawan MUFG Indonesia  |  Q2 2026  |  Pembayaran angsuran melalui potong gaji.", font_size=12, color=DARK_GREY, italic=True)
add_para(tf2, "Visi Skalabilitas", font_size=15, bold=True, color=MUFG_RED, space_before=Pt(10))
add_para(tf2, "Meluaskan platform ke seluruh klien korporat MUFG di Indonesia, dengan Samsung EPP sebagai contoh utama distribusi klien korporat.", font_size=12, color=DARK_GREY, italic=True)

# vertical divider
add_rect(sl, Inches(6.6), Inches(1.3), Inches(0.02), Inches(5.7), fill=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — PROBLEM & OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Problem & Opportunity", "Masalah & Peluang", "02")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "The Problem  |  Masalah yang Diselesaikan")
pain_points = [
    ("For Employees  |  Untuk Karyawan",
     ["Limited access to affordable consumer financing",
      "Fragmented product discovery — no single trusted platform",
      "Manual, paper-heavy loan applications",
      "High interest rates from informal lenders"],
     ["Akses terbatas ke pembiayaan konsumen terjangkau",
      "Temuan produk yang tersebar — tidak ada platform terpercaya",
      "Proses aplikasi pinjaman manual & berbasis kertas",
      "Bunga tinggi dari pemberi pinjaman informal"]),
    ("For HR / Corporate  |  Untuk HR / Korporat",
     ["No unified platform to manage employee financial benefits",
      "Payroll deduction for loans is manual and error-prone",
      "Limited visibility into employee financial wellness"],
     ["Tidak ada platform terpadu untuk manfaat keuangan karyawan",
      "Potong gaji untuk cicilan masih manual & rawan kesalahan",
      "Visibilitas terbatas terhadap kesejahteraan keuangan karyawan"]),
    ("For Providers  |  Untuk Mitra",
     ["High customer acquisition cost (CAC)",
      "Low conversion from traditional channels",
      "Repayment risk from unverified borrowers"],
     ["Biaya akuisisi pelanggan (CAC) yang tinggi",
      "Konversi rendah dari saluran tradisional",
      "Risiko gagal bayar dari peminjam yang tidak terverifikasi"]),
]
for i, (title, en_pts, bi_pts) in enumerate(pain_points):
    bx = Inches(0.3) + i * Inches(4.35)
    add_rect(sl, bx, Inches(1.25), Inches(4.1), Inches(5.7), fill=WHITE, line=MUFG_RED, line_w=Pt(1.5))
    add_rect(sl, bx, Inches(1.25), Inches(4.1), Inches(0.45), fill=MUFG_RED)
    add_text(sl, title, bx + Inches(0.05), Inches(1.3), Inches(4.0), Inches(0.4),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = sl.shapes.add_textbox(bx + Inches(0.15), Inches(1.85), Inches(3.8), Inches(2.5))
    tf = txb.text_frame; tf.word_wrap = True
    for pt in en_pts:
        add_para(tf, f"• {pt}", font_size=11, color=MUFG_DARK, space_before=Pt(3))
    txb2 = sl.shapes.add_textbox(bx + Inches(0.15), Inches(4.45), Inches(3.8), Inches(2.3))
    tf2 = txb2.text_frame; tf2.word_wrap = True
    for pt in bi_pts:
        add_para(tf2, f"• {pt}", font_size=10, color=DARK_GREY, italic=True, space_before=Pt(3))


# Opportunity slide
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Market Opportunity  |  Peluang Pasar")
stats = [
    ("IDR 850T+", "Consumer financing market in Indonesia (2025)"),
    ("68M+", "Formal sector employees — underserved for benefits"),
    ("~32%", "Employees with no access to formal credit"),
    ("3×", "Higher conversion via employer channel vs. direct-to-consumer"),
]
for i, (num, desc) in enumerate(stats):
    bx = Inches(0.3) + (i % 2) * Inches(6.55)
    by = Inches(1.3) + (i // 2) * Inches(2.95)
    add_rect(sl, bx, by, Inches(6.2), Inches(2.6), fill=WHITE, line=MUFG_RED, line_w=Pt(2))
    add_text(sl, num, bx + Inches(0.2), by + Inches(0.25), Inches(5.8), Inches(1.1),
             font_size=42, bold=True, color=MUFG_RED, align=PP_ALIGN.LEFT)
    add_text(sl, desc, bx + Inches(0.2), by + Inches(1.3), Inches(5.8), Inches(1.0),
             font_size=14, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Solution Overview", "Gambaran Solusi", "03")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Solution: MUFG Employee Benefits Portal  |  Solusi: Portal Manfaat Karyawan")

# centre value prop
add_rect(sl, Inches(0.3), Inches(1.3), Inches(12.7), Inches(1.0), fill=MUFG_DARK)
add_text(sl, "A single, white-label, payroll-integrated platform connecting MUFG corporate employees to financial products from trusted providers.",
         Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.8),
         font_size=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

pillars = [
    ("🏦  B2B2C Model", "MUFG owns the channel.\nCorporate HR administers.\nEmployees transact.",
     "Model B2B2C.\nMUFG memiliki saluran.\nHR korporat mengelola.\nKaryawan bertransaksi."),
    ("💳  Payroll Deduction\n   (Potong Gaji)", "Zero-friction repayment deducted directly from monthly salary payroll.",
     "Pembayaran tanpa gesekan langsung dipotong dari gaji bulanan."),
    ("🛒  Multi-Provider\n   Marketplace", "Adira Finance · Home Credit Indonesia · Zurich Insurance — more providers to follow.",
     "Adira Finance · Home Credit Indonesia · Zurich Insurance — mitra lebih lanjut menyusul."),
    ("📱  White-Label\n   Mobile App", "Branded for MUFG. Customisable per corporate client. API-first architecture.",
     "Bermerek MUFG. Dapat dikustomisasi per klien korporat. Arsitektur API-first."),
    ("📡  Distribution\n   Channel", "Samsung EPP and other corporate-client consumer promotions embedded in the same platform.",
     "Samsung EPP dan promosi konsumen klien korporat lainnya di platform yang sama."),
]
for i, (title, en, bi) in enumerate(pillars):
    bx = Inches(0.3) + (i % 3) * Inches(4.35) if i < 3 else Inches(0.3) + (i-3) * Inches(6.55)
    by = Inches(2.55) if i < 3 else Inches(5.05)
    w  = Inches(4.1) if i < 3 else Inches(6.2)
    add_rect(sl, bx, by, w, Inches(2.2), fill=WHITE, line=MUFG_RED, line_w=Pt(1))
    add_text(sl, title, bx + Inches(0.1), by + Inches(0.1), w - Inches(0.2), Inches(0.55),
             font_size=12, bold=True, color=MUFG_RED, align=PP_ALIGN.LEFT)
    add_text(sl, en, bx + Inches(0.1), by + Inches(0.65), w - Inches(0.2), Inches(0.7),
             font_size=11, color=MUFG_DARK, align=PP_ALIGN.LEFT)
    add_text(sl, bi, bx + Inches(0.1), by + Inches(1.35), w - Inches(0.2), Inches(0.7),
             font_size=10, color=DARK_GREY, italic=True, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — PLATFORM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Platform Architecture", "Arsitektur Platform", "04")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Platform Architecture  |  Arsitektur Platform")

# Layer diagram (text-based boxes)
layers = [
    (MUFG_RED,    "LAYER 1 — Front-End  |  Lapisan 1 — Antarmuka",
                  "White-label Mobile App (iOS / Android)  +  Employee Web Portal  |  Aplikasi Seluler & Portal Web Karyawan"),
    (ACCENT_BLUE, "LAYER 2 — MUFG Platform Core  |  Lapisan 2 — Inti Platform MUFG",
                  "Employee Identity & Auth  ·  Product Catalogue Engine  ·  Eligibility Engine  ·  Checkout & Application Flow"),
    (MUFG_DARK,   "LAYER 3 — Payroll Integration  |  Lapisan 3 — Integrasi Payroll",
                  "Payroll API / SFTP File Exchange with HR System  ·  Deduction Scheduler  ·  Reconciliation Engine"),
    (MUFG_GOLD,   "LAYER 4 — Provider APIs  |  Lapisan 4 — API Mitra",
                  "Adira Finance API  ·  Home Credit Indonesia API  ·  Zurich Insurance API  ·  Samsung EPP API"),
    (DARK_GREY,   "LAYER 5 — Data & Compliance  |  Lapisan 5 — Data & Kepatuhan",
                  "BI-OJK Compliance Layer  ·  Data Masking & Encryption  ·  Analytics Dashboard  ·  Audit Logs"),
]
for i, (col, label, detail) in enumerate(layers):
    by = Inches(1.35) + i * Inches(1.15)
    add_rect(sl, Inches(0.3), by, Inches(12.7), Inches(1.05), fill=col)
    add_text(sl, label, Inches(0.45), by + Inches(0.05), Inches(6), Inches(0.4),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(sl, detail, Inches(0.45), by + Inches(0.48), Inches(12.3), Inches(0.45),
             font_size=10, color=RGBColor(0xEE,0xEE,0xEE) if col != MUFG_GOLD else MUFG_DARK, align=PP_ALIGN.LEFT)


# User journey slide
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Employee User Journey  |  Perjalanan Pengguna Karyawan")
steps = [
    ("1", "Register\nDaftar", "Employee downloads app & authenticates via MUFG SSO\nKaryawan unduh aplikasi & autentikasi via SSO MUFG"),
    ("2", "Browse\nJelajahi", "Browse product catalogue: vehicles, devices, insurance\nJelajahi katalog: kendaraan, perangkat, asuransi"),
    ("3", "Apply\nAjukan", "Select product → fill e-form → eligibility checked in real time\nPilih produk → isi e-formulir → kelayakan dicek real-time"),
    ("4", "Approve\nDisetujui", "Provider approves in <24 hours; terms shown in-app\nMitra menyetujui <24 jam; ketentuan ditampilkan di app"),
    ("5", "Receive\nTerima", "Product delivered or policy issued; contract signed digitally\nProduk dikirim / polis diterbitkan; kontrak ditandatangani"),
    ("6", "Repay\nAngsur", "Monthly instalment auto-deducted from payroll (potong gaji)\nAngsuran bulanan dipotong otomatis dari gaji"),
]
arrow_w = Inches(1.9)
arrow_gap = Inches(0.15)
for i, (num, title, desc) in enumerate(steps):
    bx = Inches(0.25) + i * (arrow_w + arrow_gap)
    add_rect(sl, bx, Inches(1.35), arrow_w, Inches(0.55), fill=MUFG_RED)
    add_text(sl, f"STEP {num}", bx, Inches(1.38), arrow_w, Inches(0.5),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, bx, Inches(1.9), arrow_w, Inches(1.1), fill=WHITE, line=MUFG_RED, line_w=Pt(1))
    add_text(sl, title, bx + Inches(0.05), Inches(1.92), arrow_w - Inches(0.1), Inches(0.5),
             font_size=12, bold=True, color=MUFG_DARK, align=PP_ALIGN.CENTER)
    add_text(sl, desc, bx + Inches(0.05), Inches(3.1), arrow_w - Inches(0.1), Inches(3.8),
             font_size=9, color=DARK_GREY, align=PP_ALIGN.LEFT)
    if i < 5:
        add_text(sl, "▶", bx + arrow_w, Inches(1.5), arrow_gap + Inches(0.05), Inches(0.55),
                 font_size=14, bold=True, color=MUFG_RED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 5 — PROVIDER ECOSYSTEM
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Provider Ecosystem", "Ekosistem Mitra", "05")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Day-1 Provider Ecosystem  |  Ekosistem Mitra Hari Pertama")

providers = [
    ("ADIRA FINANCE", MUFG_RED,
     ["Category: Vehicle Financing (cars + motorcycles)",
      "Products: New & used car loans, motorcycle loans",
      "Tenor: 12–60 months",
      "Down payment: from 10%",
      "Rate: competitive fleet / group pricing",
      "Integration: REST API  ·  eKYC  ·  instant scoring"],
     ["Kategori: Pembiayaan Kendaraan (mobil & motor)",
      "Produk: KPM mobil baru/bekas, KPM motor",
      "Tenor: 12–60 bulan",
      "DP: mulai 10%",
      "Suku bunga: harga kompetitif grup/armada",
      "Integrasi: REST API · eKYC · scoring instan"]),
    ("HOME CREDIT INDONESIA", ACCENT_BLUE,
     ["Category: Consumer Electronics & White Goods",
      "Products: Handphones, laptops, TVs, fridges, washing machines",
      "Tenor: 3–36 months, 0% instalment promotions",
      "No collateral required",
      "Instant approval at point of sale",
      "Integration: POS API  ·  Merchant catalog sync"],
     ["Kategori: Elektronik Konsumen & Perabot Rumah Tangga",
      "Produk: HP, laptop, TV, kulkas, mesin cuci",
      "Tenor: 3–36 bulan, promo cicilan 0%",
      "Tanpa agunan",
      "Persetujuan instan di titik penjualan",
      "Integrasi: API POS · sinkronisasi katalog merchant"]),
    ("ZURICH INSURANCE", MUFG_DARK,
     ["Category: Insurance",
      "Products: Life, PA, health top-up, device protection",
      "Premium payment via payroll deduction",
      "Group rates — 20–30% below retail",
      "Digital claims via app",
      "Integration: REST API  ·  Policy issuance  ·  Claims"],
     ["Kategori: Asuransi",
      "Produk: Jiwa, PA, top-up kesehatan, proteksi perangkat",
      "Pembayaran premi via potong gaji",
      "Tarif grup — 20–30% lebih murah dari retail",
      "Klaim digital via aplikasi",
      "Integrasi: REST API · penerbitan polis · klaim"]),
]
for i, (name, col, en_pts, bi_pts) in enumerate(providers):
    bx = Inches(0.3) + i * Inches(4.35)
    add_rect(sl, bx, Inches(1.3), Inches(4.1), Inches(5.75), fill=WHITE, line=col, line_w=Pt(2))
    add_rect(sl, bx, Inches(1.3), Inches(4.1), Inches(0.55), fill=col)
    add_text(sl, name, bx + Inches(0.1), Inches(1.33), Inches(3.9), Inches(0.5),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = sl.shapes.add_textbox(bx + Inches(0.15), Inches(2.0), Inches(3.8), Inches(2.45))
    tf = txb.text_frame; tf.word_wrap = True
    for pt in en_pts:
        add_para(tf, f"• {pt}", font_size=10, color=MUFG_DARK, space_before=Pt(3))
    txb2 = sl.shapes.add_textbox(bx + Inches(0.15), Inches(4.55), Inches(3.8), Inches(2.3))
    tf2 = txb2.text_frame; tf2.word_wrap = True
    for pt in bi_pts:
        add_para(tf2, f"• {pt}", font_size=9, color=DARK_GREY, italic=True, space_before=Pt(3))


# Samsung EPP slide
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Corporate Client Distribution: Samsung EPP  |  Distribusi Klien Korporat: Samsung EPP")
add_rect(sl, Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.85), fill=MUFG_DARK)
add_text(sl, "The portal doubles as a marketing distribution channel for MUFG corporate clients' own consumer products.",
         Inches(0.5), Inches(1.38), Inches(12.3), Inches(0.6),
         font_size=13, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Portal ini juga berfungsi sebagai saluran distribusi pemasaran untuk produk konsumen milik klien korporat MUFG.",
         Inches(0.5), Inches(1.78), Inches(12.3), Inches(0.3),
         font_size=11, color=MID_GREY, italic=True, align=PP_ALIGN.CENTER)

epp_points = [
    ("What is Samsung EPP?  |  Apa itu Samsung EPP?",
     ["Employee Purchase Programme — exclusive discounts on Samsung devices for corporate employees",
      "Employees buy Galaxy phones, tablets, and wearables at 10–25% below MSRP",
      "Financed via 0% instalment through Home Credit Indonesia (already a platform partner)",
      "Samsung gets a low-CAC B2B2C channel; employees get exclusive pricing"],
     ["Employee Purchase Programme — diskon eksklusif perangkat Samsung untuk karyawan korporat",
      "Karyawan membeli Galaxy phone, tablet, dan wearable 10–25% di bawah MSRP",
      "Dibiayai cicilan 0% melalui Home Credit Indonesia (sudah menjadi mitra platform)",
      "Samsung mendapat saluran B2B2C berbiaya akuisisi rendah; karyawan mendapat harga eksklusif"]),
    ("Revenue for MUFG  |  Pendapatan untuk MUFG",
     ["Co-marketing fee from Samsung Indonesia",
      "Transaction fee per order processed through the portal",
      "Increased platform stickiness and employee engagement metrics",
      "Template for other MUFG corporate clients (e.g. automotive, FMCG brands)"],
     ["Biaya co-marketing dari Samsung Indonesia",
      "Biaya transaksi per pesanan yang diproses melalui portal",
      "Peningkatan keterikatan platform dan metrik keterlibatan karyawan",
      "Template untuk klien korporat MUFG lainnya (otomotif, FMCG, dll.)"]),
]
for i, (title, en_pts, bi_pts) in enumerate(epp_points):
    bx = Inches(0.3) + i * Inches(6.55)
    add_rect(sl, bx, Inches(2.3), Inches(6.2), Inches(4.7), fill=WHITE, line=MUFG_RED, line_w=Pt(1.5))
    add_text(sl, title, bx + Inches(0.15), Inches(2.35), Inches(5.9), Inches(0.5),
             font_size=12, bold=True, color=MUFG_RED, align=PP_ALIGN.LEFT)
    txb = sl.shapes.add_textbox(bx + Inches(0.15), Inches(2.9), Inches(5.9), Inches(2.0))
    tf = txb.text_frame; tf.word_wrap = True
    for pt in en_pts:
        add_para(tf, f"• {pt}", font_size=11, color=MUFG_DARK, space_before=Pt(4))
    txb2 = sl.shapes.add_textbox(bx + Inches(0.15), Inches(4.95), Inches(5.9), Inches(1.8))
    tf2 = txb2.text_frame; tf2.word_wrap = True
    for pt in bi_pts:
        add_para(tf2, f"• {pt}", font_size=10, color=DARK_GREY, italic=True, space_before=Pt(3))


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 6 — PILOT PROGRAM
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Pilot Program", "Program Percontohan", "06")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Pilot Program: MUFG Indonesia  |  Program Percontohan: MUFG Indonesia")

pilot_left = [
    ("Target Participants  |  Peserta Target",
     ["MUFG Bank Indonesia employees: 500–1,000 headcount",
      "Employee grades: all permanent staff",
      "HR lead: MUFG Indonesia HR Division",
      "Pilot window: Q2–Q3 2026 (6 months)"]),
    ("Success Metrics  |  Metrik Keberhasilan",
     ["≥ 30% employee registration rate within 90 days",
      "≥ 15% transaction conversion (at least 1 product taken)",
      "Zero payroll-deduction reconciliation errors",
      "NPS ≥ 40 from participating employees"]),
]
pilot_right = [
    ("Products Offered  |  Produk yang Ditawarkan",
     ["Adira: Motorcycle financing (Honda/Yamaha), Car loans",
      "Home Credit: Smartphone bundles (Samsung EPP), Appliances",
      "Zurich: Term life (Group), PA plan, Device protection"]),
    ("Operational Setup  |  Pengaturan Operasional",
     ["Payroll integration: API with MUFG HR payroll system",
      "Onboarding: In-app eKYC (KTP + selfie)",
      "Deduction schedule: monthly, aligned to payroll run",
      "Support: in-app chat + dedicated HR helpdesk"]),
]
for i, items in enumerate([pilot_left, pilot_right]):
    bx = Inches(0.3) + i * Inches(6.55)
    by = Inches(1.3)
    for j, (title, pts) in enumerate(items):
        box_y = by + j * Inches(2.9)
        add_rect(sl, bx, box_y, Inches(6.2), Inches(2.65), fill=WHITE, line=MUFG_RED, line_w=Pt(1.5))
        add_text(sl, title, bx + Inches(0.15), box_y + Inches(0.1), Inches(5.9), Inches(0.45),
                 font_size=12, bold=True, color=MUFG_RED, align=PP_ALIGN.LEFT)
        txb = sl.shapes.add_textbox(bx + Inches(0.15), box_y + Inches(0.6), Inches(5.9), Inches(1.9))
        tf = txb.text_frame; tf.word_wrap = True
        for pt in pts:
            add_para(tf, f"• {pt}", font_size=11, color=MUFG_DARK, space_before=Pt(4))


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 7 — GO-TO-MARKET
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Go-to-Market & Distribution", "Strategi Distribusi", "07")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Go-to-Market Strategy  |  Strategi Go-to-Market")

gtm = [
    ("Phase 1\nFase 1\nQ1–Q2 2026", MUFG_RED,
     ["Build & integrate platform (MVP)",
      "Onboard 3 providers (Adira, Home Credit, Zurich)",
      "Sign Samsung EPP co-marketing MOU",
      "HR system integration with MUFG Indonesia payroll"]),
    ("Phase 2\nFase 2\nQ3–Q4 2026", ACCENT_BLUE,
     ["Pilot launch: 500–1,000 MUFG Indonesia employees",
      "Soft-launch Samsung EPP on platform",
      "Collect usage data & iterate UX",
      "HR admin dashboard v1 live"]),
    ("Phase 3\nFase 3\nQ1–Q2 2027", MUFG_DARK,
     ["Expand to 5 MUFG corporate clients (target: 10,000 employees)",
      "Add 2–3 new providers (BNPL, travel insurance)",
      "Full Samsung EPP national rollout",
      "Platform-as-a-Service (PaaS) offering for MUFG clients"]),
]
for i, (phase, col, pts) in enumerate(gtm):
    bx = Inches(0.35) + i * Inches(4.35)
    add_rect(sl, bx, Inches(1.35), Inches(4.1), Inches(5.65), fill=col)
    add_text(sl, phase, bx + Inches(0.15), Inches(1.45), Inches(3.8), Inches(0.9),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = sl.shapes.add_textbox(bx + Inches(0.15), Inches(2.45), Inches(3.8), Inches(4.3))
    tf = txb.text_frame; tf.word_wrap = True
    for pt in pts:
        add_para(tf, f"• {pt}", font_size=11, color=WHITE, space_before=Pt(6))


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 8 — PRODUCT SPECIFICATION
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Product Specification", "Spesifikasi Produk", "08")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Product Specification — Mobile App  |  Spesifikasi Produk — Aplikasi Seluler")

features = [
    ("Authentication & Onboarding  |  Otentikasi",
     ["SSO via MUFG employee directory", "eKYC: KTP + selfie + liveness check",
      "Biometric login (Face ID / Fingerprint)", "Invite-code gated during pilot"]),
    ("Product Catalogue  |  Katalog Produk",
     ["Category tiles: Vehicles, Devices, Insurance, Promotions",
      "Real-time inventory & availability", "Personalised eligibility filter",
      "Samsung EPP tile (exclusive employee pricing)"]),
    ("Application Flow  |  Alur Pengajuan",
     ["In-app form: auto-populated from HR profile",
      "Real-time eligibility engine (salary vs instalment)",
      "Document upload (NPWP, KTP, salary slip)",
      "E-signature for contract execution"]),
    ("Payroll Deduction  |  Potong Gaji",
     ["Deduction schedule synced with payroll cycle",
      "Employee consent & deduction dashboard",
      "HR override & hardship pause capability",
      "Monthly statement & repayment history"]),
    ("Notifications & Support  |  Notifikasi",
     ["Push notifications: approval, deduction, promotions",
      "In-app chat (CS) + HR helpdesk escalation",
      "Claims submission (Zurich insurance)",
      "Referral & loyalty points (Phase 2)"]),
    ("Admin Dashboard  |  Dasbor Admin",
     ["HR view: headcount, deduction status, reconciliation",
      "Provider view: leads, approvals, disbursements",
      "MUFG view: revenue, GMV, platform health",
      "OJK-compliant audit trail & reporting"]),
]
for i, (title, pts) in enumerate(features):
    bx = Inches(0.3) + (i % 3) * Inches(4.35)
    by = Inches(1.3) + (i // 3) * Inches(3.05)
    add_rect(sl, bx, by, Inches(4.1), Inches(2.8), fill=WHITE, line=MUFG_RED, line_w=Pt(1))
    add_rect(sl, bx, by, Inches(4.1), Inches(0.45), fill=MUFG_RED)
    add_text(sl, title, bx + Inches(0.1), by + Inches(0.05), Inches(3.9), Inches(0.38),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txb = sl.shapes.add_textbox(bx + Inches(0.1), by + Inches(0.55), Inches(3.9), Inches(2.1))
    tf = txb.text_frame; tf.word_wrap = True
    for pt in pts:
        add_para(tf, f"• {pt}", font_size=10, color=MUFG_DARK, space_before=Pt(3))


# Tech stack slide
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Technical Architecture  |  Arsitektur Teknis")

tech_rows = [
    ("Front-End  |  Antarmuka", "React Native (iOS + Android)  ·  Next.js Web Portal  ·  MUFG Design System"),
    ("API Gateway  |  Gateway API", "Kong / AWS API Gateway  ·  OAuth 2.0 / OpenID Connect  ·  Rate limiting & WAF"),
    ("Back-End Services  |  Layanan Back-End", "Node.js microservices  ·  Eligibility Engine (Python ML)  ·  Notification Service"),
    ("Payroll Connector  |  Konektor Payroll", "REST / SFTP to HR system  ·  Deduction scheduler  ·  Reconciliation worker"),
    ("Provider Integrations  |  Integrasi Mitra", "Adira REST API  ·  Home Credit REST API  ·  Zurich Policy API  ·  Samsung EPP API"),
    ("Data & Analytics  |  Data & Analitik", "PostgreSQL + Redis  ·  Snowflake DWH  ·  Metabase dashboards  ·  Segment CDP"),
    ("Infrastructure  |  Infrastruktur", "AWS Jakarta (ap-southeast-3)  ·  Kubernetes (EKS)  ·  Terraform IaC  ·  SOC 2 / ISO 27001"),
    ("Compliance  |  Kepatuhan", "OJK POJK 6/2022 (MPMM)  ·  BI QRIS/EDC ready  ·  PDPA data residency in Indonesia"),
]
for i, (label, detail) in enumerate(tech_rows):
    col = MUFG_RED if i % 2 == 0 else MUFG_DARK
    by = Inches(1.35) + i * Inches(0.74)
    add_rect(sl, Inches(0.3), by, Inches(3.1), Inches(0.65), fill=col)
    add_text(sl, label, Inches(0.4), by + Inches(0.08), Inches(2.9), Inches(0.5),
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(sl, Inches(3.4), by, Inches(9.6), Inches(0.65), fill=WHITE, line=col, line_w=Pt(0.75))
    add_text(sl, detail, Inches(3.5), by + Inches(0.08), Inches(9.4), Inches(0.5),
             font_size=10, color=MUFG_DARK, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 9 — BUSINESS CASE & FINANCIALS
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Business Case & Financials", "Kasus Bisnis & Keuangan", "09")

# Revenue model
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Revenue Model  |  Model Pendapatan")

rev_streams = [
    ("Origination Fee  |  Biaya Originasi",
     "Charged to provider per approved application",
     "Dikenakan ke mitra per aplikasi yang disetujui",
     "0.5–1.5% of financed amount"),
    ("Platform Service Fee  |  Biaya Layanan Platform",
     "Monthly SaaS fee per corporate client on the platform",
     "Biaya SaaS bulanan per klien korporat di platform",
     "IDR 5–15M / client / month"),
    ("Co-marketing Revenue  |  Pendapatan Co-Marketing",
     "Corporate clients (e.g. Samsung) pay for feature placement",
     "Klien korporat (mis. Samsung) membayar penempatan fitur",
     "IDR 50–200M / campaign"),
    ("Interest Income  |  Pendapatan Bunga",
     "On MUFG-funded instalment products (future phase)",
     "Pada produk cicilan yang didanai MUFG (fase berikutnya)",
     "Margin 3–5% p.a."),
    ("Data & Analytics  |  Layanan Data",
     "Anonymised workforce financial analytics for providers",
     "Analitik keuangan tenaga kerja anonim untuk mitra",
     "IDR 10–50M / report"),
]
for i, (title, en, bi, range_val) in enumerate(rev_streams):
    bx = Inches(0.3) + (i % 3) * Inches(4.35) if i < 3 else Inches(0.3) + (i-3) * Inches(6.55)
    by = Inches(1.3) if i < 3 else Inches(4.4)
    w  = Inches(4.1) if i < 3 else Inches(6.2)
    add_rect(sl, bx, by, w, Inches(2.8), fill=WHITE, line=MUFG_RED, line_w=Pt(1.5))
    add_rect(sl, bx, by, w, Inches(0.45), fill=MUFG_RED)
    add_text(sl, title, bx + Inches(0.1), by + Inches(0.06), w - Inches(0.2), Inches(0.35),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(sl, en, bx + Inches(0.1), by + Inches(0.55), w - Inches(0.2), Inches(0.45),
             font_size=11, color=MUFG_DARK, align=PP_ALIGN.LEFT)
    add_text(sl, bi, bx + Inches(0.1), by + Inches(1.0), w - Inches(0.2), Inches(0.55),
             font_size=10, color=DARK_GREY, italic=True, align=PP_ALIGN.LEFT)
    add_rect(sl, bx + Inches(0.1), by + Inches(1.7), w - Inches(0.2), Inches(0.35), fill=LIGHT_GREY)
    add_text(sl, f"Est. rate: {range_val}", bx + Inches(0.15), by + Inches(1.73), w - Inches(0.3), Inches(0.3),
             font_size=10, bold=True, color=MUFG_RED, align=PP_ALIGN.LEFT)


# 3-year P&L projection
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "3-Year Financial Projection (IDR Billion)  |  Proyeksi Keuangan 3 Tahun (IDR Miliar)")

# table header
cols_w = [Inches(3.8), Inches(2.6), Inches(2.6), Inches(2.6)]
cols_x = [Inches(0.35), Inches(4.15), Inches(6.75), Inches(9.35)]
headers = ["Item  |  Keterangan", "Year 1 — Pilot\nTahun 1 (Percontohan)", "Year 2 — Scale\nTahun 2 (Skalabilitas)", "Year 3 — Growth\nTahun 3 (Pertumbuhan)"]
for j, (hdr, cx, cw) in enumerate(zip(headers, cols_x, cols_w)):
    add_rect(sl, cx, Inches(1.3), cw, Inches(0.65), fill=MUFG_RED)
    add_text(sl, hdr, cx + Inches(0.05), Inches(1.32), cw - Inches(0.1), Inches(0.6),
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# table rows
rows = [
    ("Employees on Platform  |  Karyawan di Platform", "750", "8,000", "25,000"),
    ("Corporate Clients  |  Klien Korporat", "1", "5", "15"),
    ("GMV (IDR Bn)  |  Nilai Transaksi", "7.5", "80", "250"),
    ("Origination Fee Revenue  |  Pendapatan Biaya Originasi", "0.9", "9.6", "30.0"),
    ("Platform SaaS Revenue  |  Pendapatan SaaS", "0.06", "0.36", "1.08"),
    ("Co-marketing Revenue  |  Pendapatan Co-Marketing", "0.2", "0.6", "1.5"),
    ("Total Revenue  |  Total Pendapatan", "1.16", "10.56", "32.58"),
    ("Platform OPEX  |  Biaya Operasional Platform", "(1.8)", "(3.5)", "(6.0)"),
    ("CAPEX (amortised)  |  CAPEX (diamortisasi)", "(1.2)", "(0.6)", "(0.4)"),
    ("EBITDA  |  EBITDA", "(1.84)", "6.46", "26.18"),
    ("EBITDA Margin  |  Margin EBITDA", "—", "61%", "80%"),
]
row_colors = [LIGHT_GREY, WHITE, LIGHT_GREY, WHITE, LIGHT_GREY, WHITE,
              RGBColor(0xFF,0xE5,0xE5), LIGHT_GREY, WHITE,
              RGBColor(0xC8,0x10,0x2E), LIGHT_GREY]
text_colors = [MUFG_DARK]*9 + [WHITE, MUFG_DARK]
for ri, (row_label, y1, y2, y3) in enumerate(rows):
    by = Inches(1.95) + ri * Inches(0.49)
    vals = [row_label, y1, y2, y3]
    for j, (val, cx, cw) in enumerate(zip(vals, cols_x, cols_w)):
        add_rect(sl, cx, by, cw, Inches(0.46), fill=row_colors[ri])
        bold = ri in [6, 9]
        add_text(sl, val, cx + Inches(0.05), by + Inches(0.06), cw - Inches(0.1), Inches(0.36),
                 font_size=10, bold=bold, color=text_colors[ri],
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

add_text(sl, "* Projections based on pilot conversion rate of 15%, avg ticket IDR 10M, 3-provider revenue share model. Samsung EPP co-marketing included from Year 1.",
         Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.25),
         font_size=8, color=DARK_GREY, italic=True)


# Investment required
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Investment Required & ROI  |  Investasi yang Dibutuhkan & ROI")

inv_items = [
    ("Platform Build (MVP)  |  Pembangunan Platform (MVP)", "IDR 8–10 Billion",
     "Mobile app + web portal + API integrations + payroll connector\nAplikasi mobile + portal web + integrasi API + konektor payroll"),
    ("Pilot Operations (6 months)  |  Operasional Percontohan", "IDR 1.5 Billion",
     "Staff (3 FTE), cloud infra, provider onboarding, marketing\nStaf (3 FTE), infrastruktur cloud, onboarding mitra, pemasaran"),
    ("Compliance & Licensing  |  Kepatuhan & Perizinan", "IDR 0.5 Billion",
     "OJK registration, legal, data protection audits\nPendaftaran OJK, hukum, audit perlindungan data"),
    ("Total Investment  |  Total Investasi", "IDR 10–12 Billion",
     "Estimated total for MVP + pilot phase\nEstimasi total untuk fase MVP + percontohan"),
]
for i, (title, amount, detail) in enumerate(inv_items):
    by = Inches(1.3) + i * Inches(1.4)
    is_total = i == 3
    col = MUFG_RED if is_total else MUFG_DARK
    add_rect(sl, Inches(0.3), by, Inches(12.7), Inches(1.25), fill=col if is_total else WHITE,
             line=col, line_w=Pt(1.5))
    add_text(sl, title, Inches(0.45), by + Inches(0.08), Inches(7), Inches(0.45),
             font_size=12, bold=True, color=WHITE if is_total else MUFG_DARK, align=PP_ALIGN.LEFT)
    add_text(sl, amount, Inches(9.5), by + Inches(0.05), Inches(3.3), Inches(0.55),
             font_size=18, bold=True, color=WHITE if is_total else MUFG_RED, align=PP_ALIGN.RIGHT)
    add_text(sl, detail, Inches(0.45), by + Inches(0.6), Inches(9), Inches(0.5),
             font_size=10, color=WHITE if is_total else DARK_GREY, italic=True, align=PP_ALIGN.LEFT)

# Payback period note
add_rect(sl, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.12), fill=MUFG_GOLD)
add_text(sl, "Projected payback period: ~18 months from commercial launch  ·  Break-even: end of Year 2  ·  IRR > 40% over 5 years",
         Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.3),
         font_size=10, bold=True, color=MUFG_DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 10 — ROADMAP & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
section_divider(prs, "Roadmap & Next Steps", "Peta Jalan & Langkah Selanjutnya", "10")

sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Implementation Roadmap  |  Peta Jalan Implementasi")

# Timeline bar
milestones = [
    ("Q1 2026", ["Platform design & architecture", "Provider contract negotiations",
                 "HR system API discovery", "Samsung EPP MOU signing"]),
    ("Q2 2026", ["MVP development (mobile + API)", "Payroll integration build",
                 "eKYC integration (Adira, Home Credit, Zurich)", "Internal UAT & security audit"]),
    ("Q3 2026", ["Pilot LAUNCH: MUFG Indonesia employees (500–1,000)",
                 "Samsung EPP soft-launch", "Weekly ops review cadence", "Data collection & UX iteration"]),
    ("Q4 2026", ["Pilot results analysis & board report",
                 "Corporate client 2–5 onboarding", "OJK reporting framework live", "Phase 2 planning"]),
    ("Q1 2027", ["Scale to 5 clients / 10,000 employees",
                 "Add BNPL & travel insurance providers", "Samsung EPP full national rollout",
                 "PaaS offering for MUFG clients"]),
]
for i, (qtr, tasks) in enumerate(milestones):
    bx = Inches(0.3) + i * Inches(2.6)
    add_rect(sl, bx, Inches(1.3), Inches(2.45), Inches(0.55), fill=MUFG_RED)
    add_text(sl, qtr, bx, Inches(1.33), Inches(2.45), Inches(0.5),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, bx, Inches(1.85), Inches(2.45), Inches(5.3), fill=WHITE, line=MUFG_RED, line_w=Pt(1))
    txb = sl.shapes.add_textbox(bx + Inches(0.1), Inches(1.95), Inches(2.3), Inches(5.1))
    tf = txb.text_frame; tf.word_wrap = True
    for t in tasks:
        add_para(tf, f"• {t}", font_size=10, color=MUFG_DARK, space_before=Pt(5))


# Next steps slide
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Immediate Next Steps  |  Langkah Segera")

next_steps = [
    ("1", "Executive Alignment  |  Penyelarasan Eksekutif",
     ["Present proposal to MUFG Indonesia CEO & COO",
      "Align Digital Banking and Retail divisions on ownership",
      "Confirm budget envelope for FY2026",
      "Appoint internal Product Owner"],
     "January 2026  |  Januari 2026"),
    ("2", "Provider LOIs  |  LOI Mitra",
     ["Issue Letter of Intent to Adira Finance",
      "Issue Letter of Intent to Home Credit Indonesia",
      "Issue Letter of Intent to Zurich Insurance Indonesia",
      "Begin Samsung EPP commercial term negotiation"],
     "January–February 2026"),
    ("3", "Technical Discovery  |  Penemuan Teknis",
     ["HR payroll API specification review",
      "Provider API documentation deep-dive",
      "Security & compliance assessment (OJK POJK 6/2022)",
      "Vendor RFP for platform development partner"],
     "February 2026  |  Februari 2026"),
    ("4", "Board Approval  |  Persetujuan Dewan",
     ["Prepare full business case for board approval",
      "Budget sign-off: IDR 10–12 Billion capex",
      "Legal entity structure for platform ownership",
      "Formal pilot agreement with MUFG Indonesia HR"],
     "March 2026  |  Maret 2026"),
]
for i, (num, title, tasks, deadline) in enumerate(next_steps):
    bx = Inches(0.3) + (i % 2) * Inches(6.55)
    by = Inches(1.3) + (i // 2) * Inches(2.95)
    add_rect(sl, bx, by, Inches(6.2), Inches(2.7), fill=WHITE, line=MUFG_RED, line_w=Pt(1.5))
    add_rect(sl, bx, by, Inches(0.55), Inches(2.7), fill=MUFG_RED)
    add_text(sl, num, bx, by + Inches(0.8), Inches(0.55), Inches(1.0),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx + Inches(0.65), by + Inches(0.1), Inches(5.4), Inches(0.45),
             font_size=12, bold=True, color=MUFG_RED, align=PP_ALIGN.LEFT)
    txb = sl.shapes.add_textbox(bx + Inches(0.65), by + Inches(0.6), Inches(5.4), Inches(1.6))
    tf = txb.text_frame; tf.word_wrap = True
    for t in tasks:
        add_para(tf, f"• {t}", font_size=10, color=MUFG_DARK, space_before=Pt(3))
    add_rect(sl, bx + Inches(0.65), by + Inches(2.35), Inches(5.4), Inches(0.25), fill=LIGHT_GREY)
    add_text(sl, f"Deadline: {deadline}", bx + Inches(0.7), by + Inches(2.35), Inches(5.3), Inches(0.25),
             font_size=9, bold=True, color=MUFG_RED, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=MUFG_DARK)
add_rect(sl, 0, 0, Inches(0.7), H, fill=MUFG_RED)
add_rect(sl, 0, H - Inches(1.8), W, Inches(1.8), fill=MUFG_RED)

add_text(sl, "Building the Future of", Inches(1.1), Inches(1.4), Inches(11), Inches(0.7),
         font_size=28, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)
add_text(sl, "Employee Financial Wellness", Inches(1.1), Inches(2.0), Inches(11), Inches(1.1),
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Membangun Masa Depan Kesejahteraan Keuangan Karyawan", Inches(1.1), Inches(3.0), Inches(11), Inches(0.6),
         font_size=18, bold=False, color=MUFG_GOLD, italic=True, align=PP_ALIGN.LEFT)
add_rect(sl, Inches(1.1), Inches(3.75), Inches(5), Inches(0.04), fill=MUFG_GOLD)
add_text(sl, "Powered by MUFG  ·  Adira  ·  Home Credit Indonesia  ·  Zurich Insurance  ·  Samsung EPP",
         Inches(1.1), Inches(3.9), Inches(11), Inches(0.5),
         font_size=13, color=MID_GREY, align=PP_ALIGN.LEFT)

add_text(sl, "CONFIDENTIAL  |  For internal discussion only  |  © 2026 MUFG Bank Ltd.",
         Inches(1.1), H - Inches(1.55), Inches(11), Inches(0.4),
         font_size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Strategy & Digital Banking Division  |  MUFG Indonesia",
         Inches(1.1), H - Inches(1.15), Inches(11), Inches(0.4),
         font_size=11, color=WHITE, align=PP_ALIGN.LEFT)

# ── Save ──────────────────────────────────────────────────────────────────────
OUTPUT = "/Users/msugirin/ClaudeWS/B2B2BC-Project/MUFG_Employee_Benefits_Portal_Pitch_Deck.pptx"
prs.save(OUTPUT)
print(f"✓ Saved: {OUTPUT}")
print(f"  Slides: {len(prs.slides)}")
